"""Tests for html_to_pptx.py Phase 5 template functions.

Covers _prepare_template, _find_blank_layout, _extract_theme_fonts — the three
functions that handle reference PPTX theme inheritance. Does NOT require
Playwright; tests only the python-pptx template manipulation layer.
"""

from pathlib import Path

import pytest

pptx = pytest.importorskip("pptx", reason="python-pptx not installed")
pytest.importorskip("playwright", reason="playwright not installed")
from pptx import Presentation
from pptx.util import Inches

from hippt.html_to_pptx import (
    _HEADING_SIZE_THRESHOLD_PT,
    _extract_theme_fonts,
    _find_blank_layout,
    _prepare_template,
)

_REF_DIR = Path(__file__).parent / "fixtures"
_AGENCY = str(_REF_DIR / "agency.pptx")
_BCG = str(_REF_DIR / "bcg-template.pptx")


# ── Fixtures ────────────────────────────────────────────────────────────────


@pytest.fixture()
def pptx_16_9(tmp_path):
    """Create a minimal 16:9 PPTX with one slide."""
    prs = Presentation()
    prs.slide_width = Inches(10.0)
    prs.slide_height = Inches(5.625)
    layout = prs.slide_layouts[0]
    prs.slides.add_slide(layout)
    path = tmp_path / "test_16_9.pptx"
    prs.save(str(path))
    return str(path)


# ── TestPrepareTemplate ─────────────────────────────────────────────────────


class TestPrepareTemplate:
    def test_valid_16_9_strips_slides(self, pptx_16_9):
        prs = _prepare_template(pptx_16_9)
        assert len(prs.slides) == 0

    def test_preserves_slide_masters(self, pptx_16_9):
        prs = _prepare_template(pptx_16_9)
        assert len(prs.slide_masters) >= 1

    def test_preserves_layouts(self, pptx_16_9):
        prs = _prepare_template(pptx_16_9)
        assert sum(1 for _ in prs.slide_layouts) > 0

    def test_4_3_raises_system_exit(self):
        with pytest.raises(SystemExit):
            _prepare_template(_BCG)

    def test_nonexistent_file_raises_system_exit(self):
        with pytest.raises(SystemExit):
            _prepare_template("/nonexistent/file.pptx")

    @pytest.mark.integration
    def test_agency_strips_all_slides(self):
        prs = _prepare_template(_AGENCY)
        assert len(prs.slides) == 0
        assert len(prs.slide_masters) >= 1


# ── TestFindBlankLayout ─────────────────────────────────────────────────────


class TestFindBlankLayout:
    def test_default_presentation_finds_blank(self, pptx_16_9):
        prs = _prepare_template(pptx_16_9)
        layout = _find_blank_layout(prs)
        assert layout is not None
        assert hasattr(layout, "name")

    @pytest.mark.integration
    def test_agency_returns_default_layout(self):
        prs = _prepare_template(_AGENCY)
        layout = _find_blank_layout(prs)
        assert layout.name == "DEFAULT"


# ── TestExtractThemeFonts ───────────────────────────────────────────────────


class TestExtractThemeFonts:
    def test_returns_heading_and_body_keys(self, pptx_16_9):
        prs = _prepare_template(pptx_16_9)
        fonts = _extract_theme_fonts(prs)
        assert "heading" in fonts
        assert "body" in fonts

    def test_plain_presentation_returns_dict(self):
        prs = Presentation()
        fonts = _extract_theme_fonts(prs)
        assert isinstance(fonts, dict)
        assert "heading" in fonts
        assert "body" in fonts

    @pytest.mark.integration
    def test_agency_calibri_fonts(self):
        prs = _prepare_template(_AGENCY)
        fonts = _extract_theme_fonts(prs)
        assert fonts["heading"] == "Calibri Light"
        assert fonts["body"] == "Calibri"


# ── TestHeadingSizeThreshold ────────────────────────────────────────────────


class TestHeadingSizeThreshold:
    def test_threshold_matches_design_tokens(self):
        assert _HEADING_SIZE_THRESHOLD_PT == 22


# ── TestParseCssColor — gradient fallback (E-PPTX-006) ────────────────────


class TestParseCssColorGradient:
    def test_linear_gradient_extracts_first_hex(self):
        """E-PPTX-006: linear-gradient → extract start hex as solid."""
        from pptx.dml.color import RGBColor

        from hippt.html_to_pptx import parse_css_color

        color, alpha = parse_css_color(
            "linear-gradient(135deg, #1a1a2e 0%, #16213e 100%)"
        )
        assert alpha == 1.0
        assert color == RGBColor(0x1A, 0x1A, 0x2E)

    def test_radial_gradient_extracts_first_hex(self):
        from pptx.dml.color import RGBColor

        from hippt.html_to_pptx import parse_css_color

        color, alpha = parse_css_color("radial-gradient(circle, #ff6600, #003366)")
        assert alpha == 1.0
        assert color == RGBColor(0xFF, 0x66, 0x00)

    def test_gradient_rgb_only_falls_back(self):
        """Gradient with only rgb(), no hex → black fallback with alpha > 0."""
        from hippt.html_to_pptx import parse_css_color

        _, alpha = parse_css_color(
            "linear-gradient(to right, rgb(26,26,46), rgb(22,33,62))"
        )
        assert alpha > 0


# ── TestParseCssGradient — native gradient parsing (E-PPTX-006 root fix) ───


class TestParseCssGradient:
    def test_linear_gradient_angle_and_stops(self):
        from pptx.dml.color import RGBColor

        from hippt.html_to_pptx import parse_css_gradient

        result = parse_css_gradient("linear-gradient(135deg, #1a1a2e 0%, #16213e 100%)")
        assert result is not None
        assert len(result["stops"]) == 2
        assert result["stops"][0][0] == RGBColor(0x1A, 0x1A, 0x2E)
        assert result["stops"][0][1] == 0.0
        assert result["stops"][1][1] == 1.0
        assert result["angle"] == 315.0  # (360-135+90)%360

    def test_cardinal_angles(self):
        """Verify CSS→PPTX angle conversion for all cardinal directions."""
        from hippt.html_to_pptx import parse_css_gradient

        cases = [
            ("0deg", 90.0),  # CSS bottom-to-top → PPTX 90 (bottom-to-top)
            ("90deg", 0.0),  # CSS left-to-right → PPTX 0 (left-to-right)
            ("180deg", 270.0),  # CSS top-to-bottom → PPTX 270 (top-to-bottom)
            ("270deg", 180.0),  # CSS right-to-left → PPTX 180 (right-to-left)
        ]
        for css_angle, expected_pptx in cases:
            result = parse_css_gradient(
                f"linear-gradient({css_angle}, #000000 0%, #ffffff 100%)"
            )
            assert result is not None, f"Failed to parse {css_angle}"
            assert result["angle"] == expected_pptx, (
                f"CSS {css_angle}: expected PPTX {expected_pptx}, got {result['angle']}"
            )

    def test_direction_keyword(self):
        from hippt.html_to_pptx import parse_css_gradient

        result = parse_css_gradient(
            "linear-gradient(to right, #000000 0%, #ffffff 100%)"
        )
        assert result is not None
        # CSS "to right" = 90deg → pptx angle = (360-90+90)%360 = 0
        assert result["angle"] == 0.0

    def test_rgb_stops(self):
        from hippt.html_to_pptx import parse_css_gradient

        result = parse_css_gradient(
            "linear-gradient(180deg, rgb(26,26,46) 0%, rgb(22,33,62) 100%)"
        )
        assert result is not None
        assert len(result["stops"]) == 2

    def test_transparent_stop_excluded(self):
        from hippt.html_to_pptx import parse_css_gradient

        result = parse_css_gradient(
            "linear-gradient(transparent, rgba(23,123,87,0.92))"
        )
        # transparent has alpha=0 → excluded, only 1 valid stop → returns None
        assert result is None

    def test_returns_none_for_non_gradient(self):
        from hippt.html_to_pptx import parse_css_gradient

        assert parse_css_gradient("#ff0000") is None
        assert parse_css_gradient("") is None
        assert parse_css_gradient(None) is None

    def test_radial_gradient_not_parsed(self):
        from hippt.html_to_pptx import parse_css_gradient

        assert parse_css_gradient("radial-gradient(circle, #ff6600, #003366)") is None

    def test_multi_stop_uses_first_and_last(self):
        """3+ stops collapse to first+last, preserving overall color range."""
        from pptx.dml.color import RGBColor

        from hippt.html_to_pptx import parse_css_gradient

        result = parse_css_gradient(
            "linear-gradient(90deg, #ff0000 0%, #00ff00 50%, #0000ff 100%)"
        )
        assert result is not None
        assert len(result["stops"]) == 2
        assert result["stops"][0][0] == RGBColor(0xFF, 0x00, 0x00)  # first
        assert result["stops"][1][0] == RGBColor(0x00, 0x00, 0xFF)  # last, not middle

    def test_position_clamped_to_unit_range(self):
        from hippt.html_to_pptx import parse_css_gradient

        result = parse_css_gradient("linear-gradient(90deg, #ff0000 0%, #0000ff 200%)")
        assert result is not None
        assert result["stops"][1][1] <= 1.0

    def test_gradient_after_url_found(self):
        """re.search finds gradient even when url() comes first."""
        from hippt.html_to_pptx import parse_css_gradient

        result = parse_css_gradient(
            "url(bg.png), linear-gradient(90deg, #ff0000 0%, #0000ff 100%)"
        )
        assert result is not None
        assert len(result["stops"]) == 2

    def test_native_gradient_applied_to_slide(self):
        """Integration: gradient fill with angle + positions verified."""
        from pptx.dml.color import RGBColor

        from hippt.html_to_pptx import parse_css_gradient

        prs = Presentation()
        prs.slide_width = Inches(10.0)
        prs.slide_height = Inches(5.625)
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        grad = parse_css_gradient("linear-gradient(135deg, #1a1a2e 0%, #16213e 100%)")
        fill = slide.background.fill
        fill.gradient()
        fill.gradient_angle = grad["angle"]
        for i, (color, pos) in enumerate(grad["stops"]):
            fill.gradient_stops[i].color.rgb = color
            fill.gradient_stops[i].position = pos

        assert fill.type is not None
        assert fill.gradient_stops[0].color.rgb == RGBColor(0x1A, 0x1A, 0x2E)
        assert fill.gradient_stops[1].color.rgb == RGBColor(0x16, 0x21, 0x3E)
        assert fill.gradient_angle == 315.0  # (360-135+90)%360
        assert fill.gradient_stops[0].position == 0.0
        assert fill.gradient_stops[1].position == 1.0


# ── TestTextPadding — position fidelity (E-PPTX-005 prevention) ───────────


class TestTextPadding:
    def test_padding_is_fixed_constant(self):
        """E-PPTX-005: padding must be fixed, never proportional."""
        from hippt.html_to_pptx import _TEXT_PAD_IN

        assert isinstance(_TEXT_PAD_IN, (int, float))
        assert _TEXT_PAD_IN > 0
        assert _TEXT_PAD_IN < 0.2
