#!/usr/bin/env python3
"""Analyze PPTX files to extract layout patterns for design system calibration.

Extracts per-slide metrics: canvas fill ratio, font sizes, padding, shape types,
color palette, and layout heuristics (row/column detection).

Usage:
    uv run python TOOLS/scripts/analyze_pptx.py file1.pptx [file2.pptx ...]
"""

from __future__ import annotations

import json
import logging
import statistics
import sys
from collections import Counter
from pathlib import Path
from typing import Any

log = logging.getLogger(__name__)

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ---------------------------------------------------------------------------
# Utility helpers
# ---------------------------------------------------------------------------

EMU_PER_INCH = 914400
PX_PER_INCH = 96


def emu_to_inches(emu: int) -> float:
    return round(emu / EMU_PER_INCH, 3)


def emu_to_px(emu: int) -> int:
    return round(emu / EMU_PER_INCH * PX_PER_INCH)


def hex_color(rgb_color) -> str | None:
    """Extract hex string from an RGBColor or similar object."""
    try:
        if rgb_color is None:
            return None
        r, g, b = rgb_color[0], rgb_color[1], rgb_color[2]
        return f"#{r:02X}{g:02X}{b:02X}"
    except Exception:
        return None


# ---------------------------------------------------------------------------
# Shape classification
# ---------------------------------------------------------------------------

_SHAPE_TYPE_NAMES = {
    MSO_SHAPE_TYPE.AUTO_SHAPE: "auto_shape",
    MSO_SHAPE_TYPE.CHART: "chart",
    MSO_SHAPE_TYPE.FREEFORM: "freeform",
    MSO_SHAPE_TYPE.GROUP: "group",
    MSO_SHAPE_TYPE.PICTURE: "image",
    MSO_SHAPE_TYPE.PLACEHOLDER: "placeholder",
    MSO_SHAPE_TYPE.TABLE: "table",
    MSO_SHAPE_TYPE.TEXT_BOX: "text_box",
}


def classify_shape(shape) -> str:
    """Return a human-friendly shape type string."""
    try:
        st = shape.shape_type
        return _SHAPE_TYPE_NAMES.get(
            st, str(st).split(".")[-1].lower() if st else "unknown"
        )
    except Exception:
        return "unknown"


# ---------------------------------------------------------------------------
# Core analysis functions
# ---------------------------------------------------------------------------


def analyze_slide(slide, slide_w: int, slide_h: int) -> dict[str, Any]:
    """Analyze a single slide and return metrics dict."""
    slide_area = slide_w * slide_h

    # Accumulators
    font_sizes: list[float] = []
    colors: list[str] = []
    shape_types: Counter = Counter()
    shape_positions: list[dict] = []  # for layout detection

    # Bounding box of all shapes
    bb_left = slide_w
    bb_top = slide_h
    bb_right = 0
    bb_bottom = 0
    has_shapes = False

    for shape in slide.shapes:
        shape_types[classify_shape(shape)] += 1

        # Update bounding box
        try:
            sl = shape.left or 0
            st_ = shape.top or 0
            sw = shape.width or 0
            sh = shape.height or 0
            if sw > 0 and sh > 0:
                has_shapes = True
                bb_left = min(bb_left, sl)
                bb_top = min(bb_top, st_)
                bb_right = max(bb_right, sl + sw)
                bb_bottom = max(bb_bottom, st_ + sh)
                shape_positions.append(
                    {
                        "x": sl,
                        "y": st_,
                        "w": sw,
                        "h": sh,
                        "cx": sl + sw / 2,
                        "cy": st_ + sh / 2,
                    }
                )
        except Exception as e:
            log.debug("Skipping shape bbox (missing geometry): %s", e)

        # Extract text run properties
        try:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size is not None:
                            font_sizes.append(run.font.size / Pt(1))
                        # Text color
                        try:
                            c = hex_color(run.font.color.rgb)
                            if c:
                                colors.append(c)
                        except (AttributeError, TypeError) as e:
                            log.debug("Skipping text color extraction: %s", e)
        except Exception as e:
            log.debug("Skipping text frame on shape: %s", e)

        # Fill colors from shapes
        try:
            fill = shape.fill
            if fill and fill.type is not None:
                try:
                    c = hex_color(fill.fore_color.rgb)
                    if c:
                        colors.append(c)
                except (AttributeError, TypeError) as e:
                    log.debug("Skipping fill fore_color extraction: %s", e)
        except Exception as e:
            log.debug("Skipping fill access on shape: %s", e)

    # Canvas fill ratio
    if has_shapes:
        bb_w = max(0, bb_right - bb_left)
        bb_h = max(0, bb_bottom - bb_top)
        bb_area = bb_w * bb_h
        canvas_fill = round(bb_area / slide_area, 4) if slide_area > 0 else 0
        # Clamp to 1.0 (shapes can extend beyond slide edges)
        canvas_fill = min(canvas_fill, 1.0)
    else:
        canvas_fill = 0
        bb_left = bb_top = bb_right = bb_bottom = 0

    # Padding as % of slide dimension
    if has_shapes:
        pad_left = max(0, bb_left) / slide_w if slide_w else 0
        pad_top = max(0, bb_top) / slide_h if slide_h else 0
        pad_right = max(0, slide_w - bb_right) / slide_w if slide_w else 0
        pad_bottom = max(0, slide_h - bb_bottom) / slide_h if slide_h else 0
    else:
        pad_left = pad_top = pad_right = pad_bottom = 0

    padding_pct = {
        "left": round(pad_left * 100, 1),
        "top": round(pad_top * 100, 1),
        "right": round(pad_right * 100, 1),
        "bottom": round(pad_bottom * 100, 1),
    }

    # Layout detection heuristic
    layouts = detect_layouts(shape_positions, slide_w, slide_h)

    return {
        "canvas_fill": canvas_fill,
        "font_sizes_pt": sorted(font_sizes),
        "padding_pct": padding_pct,
        "shape_types": dict(shape_types),
        "colors": colors,
        "layout": layouts,
        "shape_count": sum(shape_types.values()),
    }


def detect_layouts(positions: list[dict], slide_w: int, slide_h: int) -> list[str]:
    """Detect row/column layouts from shape positions.

    If multiple shapes share similar y-coordinates (within 5% of slide height),
    it's a row. If similar x-coordinates (within 5% of slide width), a column.
    """
    if len(positions) < 2:
        return ["single_element"] if positions else ["empty"]

    threshold_y = slide_h * 0.05
    threshold_x = slide_w * 0.05

    # Group by similar center-y (rows)
    rows: list[list[dict]] = []
    used = set()
    for i, p in enumerate(positions):
        if i in used:
            continue
        row = [p]
        used.add(i)
        for j, q in enumerate(positions):
            if j in used:
                continue
            if abs(p["cy"] - q["cy"]) < threshold_y:
                row.append(q)
                used.add(j)
        if len(row) >= 2:
            rows.append(row)

    # Group by similar center-x (columns)
    cols: list[list[dict]] = []
    used = set()
    for i, p in enumerate(positions):
        if i in used:
            continue
        col = [p]
        used.add(i)
        for j, q in enumerate(positions):
            if j in used:
                continue
            if abs(p["cx"] - q["cx"]) < threshold_x:
                col.append(q)
                used.add(j)
        if len(col) >= 2:
            cols.append(col)

    labels = []
    if rows:
        labels.append(
            f"{len(rows)} row(s) detected ({', '.join(str(len(r)) + '-across' for r in rows)})"
        )
    if cols:
        labels.append(
            f"{len(cols)} column(s) detected ({', '.join(str(len(c)) + '-down' for c in cols)})"
        )
    if not labels:
        labels.append("free_form")
    return labels


# ---------------------------------------------------------------------------
# File-level analysis
# ---------------------------------------------------------------------------


def analyze_file(path: str | Path) -> dict[str, Any]:
    """Analyze a single PPTX file and return structured results."""
    path = Path(path)
    prs = Presentation(str(path))

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    dims = {
        "emu": {"width": slide_w, "height": slide_h},
        "inches": {"width": emu_to_inches(slide_w), "height": emu_to_inches(slide_h)},
        "px_96dpi": {"width": emu_to_px(slide_w), "height": emu_to_px(slide_h)},
    }

    all_fills: list[float] = []
    all_fonts: list[float] = []
    all_colors: list[str] = []
    all_padding: list[dict] = []
    shape_type_totals: Counter = Counter()
    layout_counts: Counter = Counter()
    slide_results: list[dict] = []

    for idx, slide in enumerate(prs.slides, 1):
        result = analyze_slide(slide, slide_w, slide_h)
        result["slide_number"] = idx
        slide_results.append(result)

        all_fills.append(result["canvas_fill"])
        all_fonts.extend(result["font_sizes_pt"])
        all_colors.extend(result["colors"])
        all_padding.append(result["padding_pct"])
        shape_type_totals.update(result["shape_types"])
        for lbl in result["layout"]:
            # Normalize label for counting
            if "row" in lbl:
                layout_counts["row_layout"] += 1
            elif "column" in lbl:
                layout_counts["column_layout"] += 1
            elif lbl == "free_form":
                layout_counts["free_form"] += 1
            elif lbl == "single_element":
                layout_counts["single_element"] += 1
            elif lbl == "empty":
                layout_counts["empty"] += 1

    # Aggregate padding
    avg_padding = {}
    if all_padding:
        for key in ["left", "top", "right", "bottom"]:
            vals = [p[key] for p in all_padding]
            avg_padding[key] = round(statistics.mean(vals), 1)

    # Color palette (top 10 by frequency)
    color_freq = Counter(all_colors).most_common(10)
    palette = [{"hex": c, "count": n} for c, n in color_freq]

    # Font size stats
    font_stats = {}
    if all_fonts:
        font_stats = {
            "min": round(min(all_fonts), 1),
            "max": round(max(all_fonts), 1),
            "median": round(statistics.median(all_fonts), 1),
            "mean": round(statistics.mean(all_fonts), 1),
        }

    # Categorize fonts into title vs body heuristic
    # Title = top quartile, body = bottom 75%
    title_font_range = [0, 0]
    body_font_range = [0, 0]
    if all_fonts:
        sorted_fonts = sorted(all_fonts)
        q75_idx = int(len(sorted_fonts) * 0.75)
        body_fonts = sorted_fonts[:q75_idx] if q75_idx > 0 else sorted_fonts
        title_fonts = (
            sorted_fonts[q75_idx:] if q75_idx < len(sorted_fonts) else sorted_fonts
        )
        if title_fonts:
            title_font_range = [round(min(title_fonts), 1), round(max(title_fonts), 1)]
        if body_fonts:
            body_font_range = [round(min(body_fonts), 1), round(max(body_fonts), 1)]

    return {
        "file": path.name,
        "slide_count": len(prs.slides),
        "dimensions": dims,
        "summary": {
            "avg_canvas_fill": round(statistics.mean(all_fills), 4) if all_fills else 0,
            "median_canvas_fill": round(statistics.median(all_fills), 4)
            if all_fills
            else 0,
            "min_canvas_fill": round(min(all_fills), 4) if all_fills else 0,
            "max_canvas_fill": round(max(all_fills), 4) if all_fills else 0,
            "font_stats": font_stats,
            "title_font_pt": title_font_range,
            "body_font_pt": body_font_range,
            "avg_padding_pct": avg_padding,
            "shape_types_total": dict(shape_type_totals),
            "color_palette": palette,
            "layout_distribution": dict(layout_counts),
        },
        "slides": slide_results,
    }


def compute_cross_file_aggregates(results: list[dict]) -> dict[str, Any]:
    """Compute aggregate stats across multiple files."""
    all_fills = []
    all_fonts = []
    all_colors: Counter = Counter()
    all_padding_vals = {"left": [], "top": [], "right": [], "bottom": []}

    for r in results:
        for s in r["slides"]:
            all_fills.append(s["canvas_fill"])
            all_fonts.extend(s["font_sizes_pt"])
            all_colors.update(s["colors"])
            for key in all_padding_vals:
                all_padding_vals[key].append(s["padding_pct"][key])

    agg: dict[str, Any] = {
        "total_slides": sum(r["slide_count"] for r in results),
        "total_files": len(results),
    }

    if all_fills:
        agg["canvas_fill"] = {
            "mean": round(statistics.mean(all_fills), 4),
            "median": round(statistics.median(all_fills), 4),
            "stdev": round(statistics.stdev(all_fills), 4) if len(all_fills) > 1 else 0,
            "min": round(min(all_fills), 4),
            "max": round(max(all_fills), 4),
        }

    if all_fonts:
        agg["font_sizes_pt"] = {
            "min": round(min(all_fonts), 1),
            "max": round(max(all_fonts), 1),
            "median": round(statistics.median(all_fonts), 1),
        }

    avg_pad = {}
    for key, vals in all_padding_vals.items():
        if vals:
            avg_pad[key] = round(statistics.mean(vals), 1)
    agg["avg_padding_pct"] = avg_pad

    agg["top_colors"] = [{"hex": c, "count": n} for c, n in all_colors.most_common(15)]

    return agg


# ---------------------------------------------------------------------------
# Dominant pattern detection
# ---------------------------------------------------------------------------


def detect_dominant_patterns(result: dict) -> list[str]:
    """Generate human-readable pattern descriptions from a file's analysis."""
    patterns = []
    summary = result["summary"]

    # Canvas fill
    avg_fill = summary["avg_canvas_fill"]
    if avg_fill > 0.8:
        patterns.append(
            f"Full-bleed dominant: avg canvas fill {avg_fill:.0%}, content extends to edges"
        )
    elif avg_fill > 0.6:
        patterns.append(
            f"High fill: avg canvas fill {avg_fill:.0%}, generous content with moderate margins"
        )
    elif avg_fill > 0.4:
        patterns.append(
            f"Balanced layout: avg canvas fill {avg_fill:.0%}, clear breathing room"
        )
    else:
        patterns.append(
            f"Minimal/sparse layout: avg canvas fill {avg_fill:.0%}, large whitespace margins"
        )

    # Layout distribution
    ld = summary.get("layout_distribution", {})
    total_layouts = sum(ld.values()) or 1
    if ld.get("row_layout", 0) / total_layouts > 0.3:
        patterns.append("Row-based compositions frequent (side-by-side elements)")
    if ld.get("column_layout", 0) / total_layouts > 0.3:
        patterns.append("Column-based stacking frequent (vertical element grouping)")
    if ld.get("free_form", 0) / total_layouts > 0.4:
        patterns.append("Free-form positioning dominant (non-grid layouts)")

    # Font range
    font_stats = summary.get("font_stats", {})
    if font_stats:
        rng = font_stats["max"] - font_stats["min"]
        if rng > 30:
            patterns.append(
                f"High typographic contrast: {font_stats['min']}pt to {font_stats['max']}pt (range {rng:.0f}pt)"
            )
        elif rng > 15:
            patterns.append(
                f"Moderate type scale: {font_stats['min']}pt to {font_stats['max']}pt"
            )
        else:
            patterns.append(
                f"Tight type scale: {font_stats['min']}pt to {font_stats['max']}pt (range {rng:.0f}pt)"
            )

    # Shape composition
    st = summary.get("shape_types_total", {})
    total_shapes = sum(st.values()) or 1
    img_pct = st.get("image", 0) / total_shapes
    if img_pct > 0.3:
        patterns.append(
            f"Image-heavy: {st.get('image', 0)} images ({img_pct:.0%} of shapes)"
        )
    if st.get("chart", 0) > 0:
        patterns.append(f"Data visualization present: {st['chart']} chart(s)")
    if st.get("table", 0) > 0:
        patterns.append(f"Tabular data present: {st['table']} table(s)")

    # Padding symmetry
    pad = summary.get("avg_padding_pct", {})
    if pad:
        h_sym = abs(pad.get("left", 0) - pad.get("right", 0))
        v_sym = abs(pad.get("top", 0) - pad.get("bottom", 0))
        if h_sym < 2 and v_sym < 2:
            patterns.append("Symmetric padding (even margins on all sides)")
        elif h_sym < 2:
            patterns.append("Horizontally symmetric padding, vertical asymmetry")

    return patterns


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------


def main(paths: list[str]) -> dict[str, Any]:
    """Analyze multiple PPTX files and return combined results."""
    results = []
    for p in paths:
        print(f"Analyzing {Path(p).name}...", file=sys.stderr)
        r = analyze_file(p)
        r["patterns"] = detect_dominant_patterns(r)
        results.append(r)

    cross = compute_cross_file_aggregates(results)

    output = {
        "files": results,
        "cross_file_aggregates": cross,
    }
    return output


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: analyze_pptx.py <file1.pptx> [file2.pptx ...]", file=sys.stderr)
        sys.exit(1)

    result = main(sys.argv[1:])

    # Output JSON (exclude per-slide detail for cleaner output; use --verbose for full)
    if "--verbose" in sys.argv:
        print(json.dumps(result, indent=2))
    else:
        # Compact: drop per-slide details from output
        compact = {
            "files": [],
            "cross_file_aggregates": result["cross_file_aggregates"],
        }
        for f in result["files"]:
            compact_file = {k: v for k, v in f.items() if k != "slides"}
            compact["files"].append(compact_file)
        print(json.dumps(compact, indent=2))
