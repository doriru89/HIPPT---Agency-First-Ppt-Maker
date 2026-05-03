"""
pptx_fidelity.py — Visual fidelity comparison: HTML screenshots vs PPTX structure.

Captures HTML slides via Playwright, extracts PPTX element data via pptx_to_layout,
compares per-slide, and generates a side-by-side comparison HTML + findings JSON.
"""

import json
import logging
import threading
from dataclasses import dataclass, field
from datetime import datetime
from http.server import HTTPServer, SimpleHTTPRequestHandler
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

log = logging.getLogger(__name__)

VIEWPORT_W = 960
VIEWPORT_H = 540


@dataclass
class CaptureManifest:
    """Instructions for MCP Playwright to capture Google Slides screenshots."""

    url: str
    viewport: tuple[int, int]
    slide_count: int
    out_dir: Path
    steps: list[dict] = field(default_factory=list)

    def expected_paths(self) -> list[Path]:
        return [self.out_dir / f"pptx-s{i}.png" for i in range(1, self.slide_count + 1)]


@dataclass
class SlideFindings:
    """Per-slide comparison results."""

    slide_idx: int
    position_diffs: list[dict] = field(default_factory=list)
    font_diffs: list[dict] = field(default_factory=list)
    missing_content: list[str] = field(default_factory=list)
    extra_elements: list[str] = field(default_factory=list)
    element_count_html: int = 0
    element_count_pptx: int = 0

    def severity(self) -> str:
        max_pos = max((d["delta_in"] for d in self.position_diffs), default=0)
        if max_pos > 1.0 or self.missing_content:
            return "blocking"
        if max_pos > 0.3:
            return "degrading"
        return "cosmetic"


def capture_html_slides(html_path: Path, out_dir: Path) -> list[Path]:
    """Screenshot each .slide element at 960×540 via local HTTP + Playwright."""
    from playwright.sync_api import sync_playwright

    out_dir.mkdir(parents=True, exist_ok=True)
    html_dir = str(html_path.parent)
    html_name = html_path.name

    class _QuietHandler(SimpleHTTPRequestHandler):
        def __init__(self, *args, **kwargs):
            super().__init__(*args, directory=html_dir, **kwargs)

        def log_message(self, *args):
            pass

    server = HTTPServer(("127.0.0.1", 0), _QuietHandler)
    port = server.server_address[1]
    threading.Thread(target=server.serve_forever, daemon=True).start()

    paths = []
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(
                viewport={"width": VIEWPORT_W, "height": VIEWPORT_H},
                device_scale_factor=1,
            )
            page.goto(f"http://127.0.0.1:{port}/{html_name}", wait_until="networkidle")
            page.wait_for_timeout(500)

            slides = page.query_selector_all(".slide")
            if not slides:
                log.error("No .slide elements found in %s", html_path)
                return []

            for i, slide_el in enumerate(slides, 1):
                out_file = out_dir / f"html-s{i}.png"
                slide_el.screenshot(path=str(out_file))
                paths.append(out_file)
                log.info("Captured HTML slide %d → %s", i, out_file.name)

            browser.close()
    finally:
        server.shutdown()

    return paths


def gslides_capture_manifest(
    file_id: str, slide_count: int, out_dir: Path
) -> CaptureManifest:
    """Build MCP Playwright instructions for Google Slides screenshot capture."""
    url = f"https://docs.google.com/presentation/d/{file_id}/present"
    manifest = CaptureManifest(
        url=url,
        viewport=(VIEWPORT_W, VIEWPORT_H),
        slide_count=slide_count,
        out_dir=out_dir,
    )
    for i in range(1, slide_count + 1):
        manifest.steps.append(
            {
                "action": "screenshot" if i == 1 else "press_right_then_screenshot",
                "wait_ms": 1500,
                "save_path": str(out_dir / f"pptx-s{i}.png"),
            }
        )
    return manifest


def extract_pptx_structure(pptx_path: Path) -> list[list[dict]]:
    """Extract per-slide element data from PPTX using pptx_to_layout pattern."""
    prs = Presentation(str(pptx_path))
    slide_w_emu = prs.slide_width
    slide_h_emu = prs.slide_height

    all_slides = []
    for slide in prs.slides:
        elements = []
        for shape in slide.shapes:
            left = shape.left or 0
            top = shape.top or 0
            width = shape.width or 0
            height = shape.height or 0

            x_in = Emu(left).inches
            y_in = Emu(top).inches
            w_in = Emu(width).inches
            h_in = Emu(height).inches

            text = ""
            font_pt = 0.0
            if shape.has_text_frame:
                text = shape.text_frame.text
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            font_pt = max(font_pt, run.font.size.pt)

            elements.append(
                {
                    "name": shape.name,
                    "x_in": round(x_in, 3),
                    "y_in": round(y_in, 3),
                    "w_in": round(w_in, 3),
                    "h_in": round(h_in, 3),
                    "text": text[:100],
                    "font_pt": round(font_pt, 1),
                }
            )
        all_slides.append(elements)
    return all_slides


def compare_slide(
    sidecar_elements: list[dict],
    pptx_elements: list[dict],
    slide_idx: int,
) -> SlideFindings:
    """Compare sidecar JSON elements against PPTX extracted elements."""
    findings = SlideFindings(
        slide_idx=slide_idx,
        element_count_html=len(sidecar_elements),
        element_count_pptx=len(pptx_elements),
    )

    pptx_by_name = {e["name"]: e for e in pptx_elements}
    pptx_by_text = {}
    for e in pptx_elements:
        if e["text"]:
            pptx_by_text.setdefault(e["text"][:50], []).append(e)

    matched_pptx = set()

    for i, sidecar_el in enumerate(sidecar_elements):
        elem_type = sidecar_el.get("type", "text")
        expected_name = f"s{slide_idx}-{elem_type}-{i}"
        content = str(sidecar_el.get("content", sidecar_el.get("value", "")))[:50]

        pptx_el = pptx_by_name.get(expected_name)
        if not pptx_el and content:
            candidates = pptx_by_text.get(content[:50], [])
            for c in candidates:
                if c["name"] not in matched_pptx:
                    pptx_el = c
                    break

        if not pptx_el:
            has_pos = "x" in sidecar_el and "y" in sidecar_el
            if has_pos:
                s_x = sidecar_el["x"]
                s_y = sidecar_el["y"]
                best_dist = float("inf")
                best_candidate = None
                for candidate in pptx_elements:
                    if candidate["name"] in matched_pptx:
                        continue
                    c_x = candidate["x_in"]
                    c_y = candidate["y_in"]
                    dist = ((s_x - c_x) ** 2 + (s_y - c_y) ** 2) ** 0.5
                    if dist < best_dist:
                        best_dist = dist
                        best_candidate = candidate
                if best_candidate and best_dist <= 0.5:
                    pptx_el = best_candidate

        if not pptx_el:
            if content:
                findings.missing_content.append(content[:60])
            continue

        matched_pptx.add(pptx_el["name"])

        has_explicit_pos = "x" in sidecar_el and "y" in sidecar_el
        if has_explicit_pos:
            s_x = sidecar_el["x"]
            s_y = sidecar_el["y"]
            p_x = pptx_el["x_in"]
            p_y = pptx_el["y_in"]
            delta = ((s_x - p_x) ** 2 + (s_y - p_y) ** 2) ** 0.5

            if delta > 0.1:
                findings.position_diffs.append(
                    {
                        "element": expected_name,
                        "content": content[:30],
                        "sidecar_pos": (s_x, s_y),
                        "pptx_pos": (p_x, p_y),
                        "delta_in": round(delta, 2),
                    }
                )

        s_font = sidecar_el.get("font_size", 0)
        p_font = pptx_el["font_pt"]
        if s_font and p_font and abs(s_font - p_font) > 2:
            findings.font_diffs.append(
                {
                    "element": expected_name,
                    "content": content[:30],
                    "sidecar_pt": s_font,
                    "pptx_pt": p_font,
                    "delta_pt": round(abs(s_font - p_font), 1),
                }
            )

    for e in pptx_elements:
        if e["name"] not in matched_pptx and e["text"]:
            findings.extra_elements.append(f"{e['name']}: {e['text'][:40]}")

    return findings


def generate_comparison_html(
    html_shots: list[Path],
    pptx_shots: list[Path],
    findings: list[SlideFindings],
    out_path: Path,
) -> Path:
    """Generate 2-up comparison HTML + findings.json."""
    out_path.parent.mkdir(parents=True, exist_ok=True)

    n_slides = max(len(html_shots), len(pptx_shots))

    slides_html = []
    for i in range(n_slides):
        html_img = html_shots[i].name if i < len(html_shots) else ""
        pptx_img = pptx_shots[i].name if i < len(pptx_shots) else ""
        finding = findings[i] if i < len(findings) else None

        severity = finding.severity() if finding else "unknown"
        severity_color = {
            "blocking": "#e74c3c",
            "degrading": "#f39c12",
            "cosmetic": "#27ae60",
        }.get(severity, "#95a5a6")

        findings_items = ""
        if finding:
            for pd in finding.position_diffs[:5]:
                findings_items += (
                    f'<li class="pos-diff">Position: {pd["content"]} '
                    f"delta={pd['delta_in']}in</li>\n"
                )
            for fd in finding.font_diffs[:5]:
                findings_items += (
                    f'<li class="font-diff">Font: {fd["content"]} '
                    f"{fd['sidecar_pt']}pt→{fd['pptx_pt']}pt</li>\n"
                )
            for mc in finding.missing_content[:3]:
                findings_items += f'<li class="missing">Missing: {mc}</li>\n'

        slides_html.append(f"""
        <div class="slide-pair">
            <h3>Slide {i + 1}
                <span class="severity" style="background:{severity_color}">{severity}</span>
            </h3>
            <div class="images">
                <div class="img-col">
                    <label>HTML (ground truth)</label>
                    {"<img src='" + html_img + "'>" if html_img else "<p>No screenshot</p>"}
                </div>
                <div class="img-col">
                    <label>PPTX (Google Slides)</label>
                    {"<img src='" + pptx_img + "'>" if pptx_img else "<p>No screenshot</p>"}
                </div>
            </div>
            <ul class="findings">{findings_items}</ul>
        </div>
        """)

    html_content = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8">
<title>PPTX Fidelity Comparison</title>
<style>
body {{ font-family: -apple-system, sans-serif; background: #1a1a2e; color: #eee;
       max-width: 1920px; margin: 0 auto; padding: 24px; }}
h1 {{ text-align: center; }}
.slide-pair {{ background: #16213e; border-radius: 8px; padding: 16px;
              margin-bottom: 24px; }}
.slide-pair h3 {{ margin: 0 0 12px; display: flex; align-items: center; gap: 12px; }}
.severity {{ font-size: 11px; padding: 2px 8px; border-radius: 4px; color: #fff; }}
.images {{ display: grid; grid-template-columns: 1fr 1fr; gap: 12px; }}
.img-col {{ text-align: center; }}
.img-col label {{ display: block; font-size: 12px; color: #7a96ad; margin-bottom: 4px; }}
.img-col img {{ width: 100%; border: 1px solid #2a3a5e; border-radius: 4px; }}
.findings {{ font-size: 12px; margin-top: 8px; padding-left: 20px; }}
.pos-diff {{ color: #f39c12; }}
.font-diff {{ color: #3498db; }}
.missing {{ color: #e74c3c; }}
.meta {{ text-align: center; color: #7a96ad; font-size: 12px; margin-bottom: 20px; }}
</style></head><body>
<h1>PPTX Fidelity Comparison</h1>
<p class="meta">Generated {datetime.now().strftime("%Y-%m-%d %H:%M")} | {n_slides} slides</p>
{"".join(slides_html)}
</body></html>"""

    out_path.write_text(html_content)

    findings_json = {
        "generated": datetime.now().isoformat(),
        "slide_count": n_slides,
        "slides": [
            {
                "idx": f.slide_idx,
                "severity": f.severity(),
                "position_diffs": f.position_diffs,
                "font_diffs": f.font_diffs,
                "missing_content": f.missing_content,
                "extra_elements": f.extra_elements,
                "element_counts": {
                    "html": f.element_count_html,
                    "pptx": f.element_count_pptx,
                },
            }
            for f in findings
        ],
    }
    json_path = out_path.with_name("findings.json")
    json_path.write_text(json.dumps(findings_json, indent=2))

    return out_path
