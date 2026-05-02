#!/usr/bin/env python3
"""Draft PPTX engine — semantic sidecar JSON, fast, no browser.

Reads structured slide data from a sidecar JSON file (emitted by /presentation
--mode deck) and design tokens from a YAML file, then maps each element to
python-pptx calls. 15 specialized handlers encode slide semantics (radial,
era_card, stat_hero, etc.).

Counterpart: html_to_pptx.py = final PPTX engine (DOM extraction, Playwright,
pixel-accurate). Both coexist permanently.

Usage:
    uv run python TOOLS/scripts/slides_to_pptx.py <slides.json> [--tokens <ref.yaml>] [--out <output.pptx>] [--engine python-pptx|pptxgenjs]
"""

from __future__ import annotations

import json
import logging
import sys
from datetime import date
from pathlib import Path


import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.presentation import Presentation as PresentationClass
from pptx.util import Inches, Pt

from hippt.design_tokens import (
    ALIGN_CSS_TO_PPTX,
    COL_WIDTH_IN,
    DEFAULT_FONT,
    MARGIN_IN,
    SAFE_FONTS as _DT_SAFE_FONTS,
    SLIDE_HEIGHT_IN as SLIDE_HEIGHT,
    SLIDE_WIDTH_IN as SLIDE_WIDTH,
    hex_to_rgb as _dt_hex_to_rgb,
)

log = logging.getLogger(__name__)

SAFE_FONTS = set(_DT_SAFE_FONTS)

ALIGN_MAP = {k: getattr(PP_ALIGN, v) for k, v in ALIGN_CSS_TO_PPTX.items()}


def hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert '#RRGGBB' or 'RRGGBB' to RGBColor. E-PPTX-001 compliant."""
    r, g, b = _dt_hex_to_rgb(hex_str)
    return RGBColor(r, g, b)


class TokenResolver:
    """Resolves color/font references from design tokens YAML."""

    def __init__(self, tokens: dict | None = None):
        self.palette: dict[str, str] = {}
        self.fonts: dict[str, str] = {}
        self.fallback_map: dict[str, str] = {}

        if not tokens:
            return

        palette = tokens.get("palette") or tokens.get("colors") or []
        if isinstance(palette, dict):
            for role, hex_val in palette.items():
                if role and hex_val and role not in self.palette:
                    self.palette[role] = hex_val
        else:
            for entry in palette:
                role = entry.get("role", "")
                hex_val = entry.get("hex", "")
                if role and hex_val and role not in self.palette:
                    self.palette[role] = hex_val

        typo = tokens.get("typography") or tokens.get("fonts") or []
        if isinstance(typo, dict):
            for role, spec in typo.items():
                if isinstance(spec, dict):
                    family = spec.get("family", "")
                elif isinstance(spec, str):
                    family = spec
                else:
                    family = ""
                if role and family:
                    self.fonts[role] = family
        else:
            for entry in typo:
                role = entry.get("role", "")
                family = entry.get("family", "")
                if role and family:
                    self.fonts[role] = family

    def resolve_color(self, color_ref: str | None) -> RGBColor | None:
        """Resolve a color reference: hex string, role name, or None."""
        if not color_ref:
            return None
        if color_ref in self.palette:
            return hex_to_rgb(self.palette[color_ref])
        if color_ref.startswith("#") or len(color_ref) == 6:
            return hex_to_rgb(color_ref)
        return hex_to_rgb("1E293B")

    def resolve_font(self, font_ref: str | None) -> str:
        """Resolve font with safe fallback. E-PPTX-003 compliant."""
        if not font_ref:
            return self.fonts.get("body", DEFAULT_FONT)
        if font_ref in self.fonts:
            resolved = self.fonts[font_ref]
        else:
            resolved = font_ref

        if resolved in SAFE_FONTS:
            return resolved

        if resolved in self.fallback_map:
            return self.fallback_map[resolved]

        fallback = DEFAULT_FONT
        self.fallback_map[resolved] = fallback
        log.warning(
            "Font '%s' not in safe list, falling back to '%s'", resolved, fallback
        )
        return fallback


def set_background(slide, bg: dict | str | None, resolver: TokenResolver):
    """Set slide background from background spec."""
    if not bg:
        return
    if isinstance(bg, str):
        bg = {"type": "solid", "color": bg}
    bg_type = bg.get("type", "solid")
    color = bg.get("color")

    if bg_type == "solid" and color:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = resolver.resolve_color(color)
    elif bg_type == "gradient":
        fill = slide.background.fill
        fill.solid()
        color_val = bg.get("stops", [{}])[0].get("color", color)
        if color_val:
            fill.fore_color.rgb = resolver.resolve_color(color_val)


def add_textbox(
    slide, elem: dict, resolver: TokenResolver, slide_text_color: str | None = None
):
    """Add a text element as an editable textbox."""
    role = elem.get("role", "")

    # Column roles: position at half-width instead of full-width
    if role == "column_left":
        return _add_column_textbox(
            slide, elem, resolver, side="left", slide_text_color=slide_text_color
        )
    if role == "column_right":
        return _add_column_textbox(
            slide, elem, resolver, side="right", slide_text_color=slide_text_color
        )

    x = elem.get("x", 0)
    y = elem.get("y", 0)
    w = elem.get("w", 5)
    h = elem.get("h", 1)

    txbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txbox.text_frame
    tf.word_wrap = elem.get("word_wrap", True)

    content = elem.get("content", "")
    runs = elem.get("runs")
    default_color = slide_text_color or elem.get("color")

    if runs:
        for i, run_data in enumerate(runs):
            if i == 0 and len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.alignment = ALIGN_MAP.get(
                run_data.get("align", elem.get("align", "left")), PP_ALIGN.LEFT
            )
            p.space_after = Pt(run_data.get("spacing_after", 0))

            run = p.add_run()
            run.text = run_data.get("text", "")
            run.font.name = resolver.resolve_font(run_data.get("font"))
            run.font.size = Pt(run_data.get("font_size", 14))
            run.font.bold = run_data.get("bold", False)
            run.font.italic = run_data.get("italic", False)
            color = run_data.get("color") or default_color
            if color:
                run.font.color.rgb = resolver.resolve_color(color)
    elif content:
        paragraphs = content.split("\n") if "\n" in content else [content]
        for i, para_text in enumerate(paragraphs):
            if i == 0 and len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.alignment = ALIGN_MAP.get(elem.get("align", "left"), PP_ALIGN.LEFT)
            p.space_after = Pt(elem.get("spacing_after", 0))

            run = p.add_run()
            run.text = para_text
            run.font.name = resolver.resolve_font(elem.get("font"))
            run.font.size = Pt(elem.get("font_size", 14))
            run.font.bold = elem.get("bold", False)
            run.font.italic = elem.get("italic", False)
            color = default_color
            if color:
                run.font.color.rgb = resolver.resolve_color(color)


def _add_column_textbox(
    slide,
    elem: dict,
    resolver: TokenResolver,
    side: str,
    slide_text_color: str | None = None,
):
    """Render column_left/column_right editorial text with heading+content+highlight+detail."""
    half_w = (SLIDE_WIDTH - 2 * MARGIN - 0.3) / 2
    if side == "left":
        x = MARGIN
    else:
        x = MARGIN + half_w + 0.3
    y = elem.get("y", 0)
    h = elem.get("h", 3.0)
    color = slide_text_color or "1E293B"

    txbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(half_w), Inches(h))
    tf = txbox.text_frame
    tf.word_wrap = True

    heading = elem.get("heading", "")
    if heading:
        _add_simple_run(tf, heading, "display", 13, resolver, bold=True, color=color)

    content = elem.get("content", "")
    if content:
        _add_simple_run(tf, content, None, 9, resolver, color="475569")

    highlight = elem.get("highlight", "")
    if highlight:
        _add_simple_run(tf, highlight, None, 9, resolver, bold=True, color="00A896")

    detail = elem.get("detail", "")
    if detail:
        _add_simple_run(tf, detail, None, 8, resolver, color="64748B")


def add_shape(slide, elem: dict, resolver: TokenResolver):
    """Add a shape (rectangle, oval, etc.)."""
    x = elem.get("x", 0)
    y = elem.get("y", 0)
    w = elem.get("w", 1)
    h = elem.get("h", 1)

    shape_type = elem.get("shape", "rectangle")
    # python-pptx auto shape type IDs
    type_map = {"rectangle": 1, "oval": 9, "rounded_rectangle": 5, "triangle": 7}
    mso_type = type_map.get(shape_type, 1)

    shape = slide.shapes.add_shape(mso_type, Inches(x), Inches(y), Inches(w), Inches(h))

    fill_color = elem.get("fill")
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = resolver.resolve_color(fill_color)
    else:
        shape.fill.background()

    line_color = elem.get("line_color")
    if line_color:
        shape.line.color.rgb = resolver.resolve_color(line_color)
        shape.line.width = Pt(elem.get("line_width", 0.5))
    else:
        shape.line.fill.background()

    return shape


def add_image(slide, elem: dict):
    """Add an image from file path."""
    x = elem.get("x", 0)
    y = elem.get("y", 0)
    w = elem.get("w")
    h = elem.get("h")
    path = elem.get("path", "") or elem.get("src", "")

    if not Path(path).exists():
        log.warning("Image not found: %s", path)
        return

    kwargs = {"image_file": str(path), "left": Inches(x), "top": Inches(y)}
    if w:
        kwargs["width"] = Inches(w)
    if h:
        kwargs["height"] = Inches(h)

    slide.shapes.add_picture(**kwargs)


def add_table(slide, elem: dict, resolver: TokenResolver):
    """Add a table with headers and rows."""
    headers = elem.get("headers", [])
    rows_data = elem.get("rows", [])
    x = elem.get("x", 0.5)
    y = elem.get("y", 1)
    w = elem.get("w", 9)
    if not headers and not rows_data:
        return

    n_rows = len(rows_data) + (1 if headers else 0)
    n_cols = len(headers) if headers else (len(rows_data[0]) if rows_data else 1)
    has_multiline = any(
        "\n" in str(cell)
        for row in rows_data
        for cell in (row if isinstance(row, list) else [row])
    )
    row_h = 0.55 if has_multiline else 0.4
    h = elem.get("h", max(2.5, n_rows * row_h))

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tbl = table_shape.table

    font_name = resolver.resolve_font(None)
    header_color = resolver.resolve_color(elem.get("header_color", "1E293B"))
    body_color = resolver.resolve_color(elem.get("body_color", "475569"))

    if headers:
        for ci, h_text in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.text = str(h_text)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = font_name
                    r.font.size = Pt(elem.get("header_font_size", 8))
                    r.font.bold = True
                    r.font.color.rgb = header_color

    offset = 1 if headers else 0
    for ri, row in enumerate(rows_data):
        for ci, val in enumerate(row):
            if ci >= n_cols:
                break
            cell = tbl.cell(ri + offset, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = font_name
                    r.font.size = Pt(elem.get("body_font_size", 8))
                    r.font.color.rgb = body_color


MARGIN = MARGIN_IN
COL_WIDTH = COL_WIDTH_IN


def _grid_rect(
    elem: dict, y_cursor: float | None = None
) -> tuple[float, float, float, float]:
    """Derive x, y, w, h from gridColumn or explicit values.

    If y_cursor is provided and the element has no explicit 'y', uses the cursor
    position instead of the hardcoded default. This enables sequential stacking.
    """
    if "x" in elem:
        return elem["x"], elem.get("y", 1.0), elem.get("w", 4), elem.get("h", 3)
    gc = elem.get("gridColumn", "1 / -1")
    parts = [p.strip() for p in gc.split("/")]
    start = int(parts[0]) - 1
    end = 12 if parts[1].strip() == "-1" else int(parts[1]) - 1
    x = MARGIN + start * COL_WIDTH
    w = (end - start) * COL_WIDTH
    default_y = y_cursor if (y_cursor is not None and "y" not in elem) else 1.2
    return x, elem.get("y", default_y), w, elem.get("h", 3.5)


def _add_simple_run(
    tf,
    text: str,
    font_name: str | None,
    size: int,
    resolver: TokenResolver,
    bold: bool = False,
    color: str | None = None,
    align=PP_ALIGN.LEFT,
):
    """Append a paragraph with a single run to a text frame."""
    if tf.paragraphs and tf.paragraphs[0].text == "" and len(tf.paragraphs) == 1:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = resolver.resolve_font(font_name)
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = resolver.resolve_color(color)
    return p


def add_card(slide, elem: dict, resolver: TokenResolver):
    """Card with accent left border, label, content, optional detail/evidence."""
    x, y, w, h = _grid_rect(elem)
    accent = elem.get("accent", "#2563EB")

    # Accent left border
    border = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.06), Inches(h))
    border.fill.solid()
    border.fill.fore_color.rgb = resolver.resolve_color(accent)
    border.line.fill.background()

    # Background rectangle
    bg = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb("F8FAFC")
    bg.line.fill.background()

    # Title
    txbox = slide.shapes.add_textbox(
        Inches(x + 0.15), Inches(y + 0.1), Inches(w - 0.3), Inches(0.4)
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    _add_simple_run(
        tf, elem.get("label", ""), None, 11, resolver, bold=True, color="1E293B"
    )

    # Body content
    content = elem.get("content", "")
    if content:
        body = slide.shapes.add_textbox(
            Inches(x + 0.15), Inches(y + 0.5), Inches(w - 0.3), Inches(h - 1.0)
        )
        btf = body.text_frame
        btf.word_wrap = True
        _add_simple_run(btf, content, None, 9, resolver, color="475569")

    # Detail
    detail = elem.get("detail", "")
    if detail:
        dtx = slide.shapes.add_textbox(
            Inches(x + 0.15), Inches(y + h - 0.7), Inches(w - 0.3), Inches(0.3)
        )
        dtf = dtx.text_frame
        dtf.word_wrap = True
        _add_simple_run(dtf, detail, None, 8, resolver, color="64748B")

    # Evidence
    evidence = elem.get("evidence", "")
    if evidence:
        etx = slide.shapes.add_textbox(
            Inches(x + 0.15), Inches(y + h - 0.35), Inches(w - 0.3), Inches(0.3)
        )
        etf = etx.text_frame
        etf.word_wrap = True
        _add_simple_run(etf, evidence, None, 7, resolver, color="94A3B8")


def add_stat_card(slide, elem: dict, resolver: TokenResolver):
    """Stat card: large value + smaller label."""
    x, y, w, h = _grid_rect(elem)
    accent = elem.get("accent", "#2563EB")

    bg = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb("F8FAFC")
    bg.line.fill.background()

    val_box = slide.shapes.add_textbox(
        Inches(x + 0.1), Inches(y + 0.15), Inches(w - 0.2), Inches(h * 0.55)
    )
    vtf = val_box.text_frame
    vtf.word_wrap = True
    _add_simple_run(
        vtf,
        str(elem.get("value", "")),
        None,
        22,
        resolver,
        bold=True,
        color=accent,
        align=PP_ALIGN.CENTER,
    )

    lbl_box = slide.shapes.add_textbox(
        Inches(x + 0.1), Inches(y + h * 0.55), Inches(w - 0.2), Inches(h * 0.35)
    )
    ltf = lbl_box.text_frame
    ltf.word_wrap = True
    _add_simple_run(
        ltf,
        elem.get("label", ""),
        None,
        9,
        resolver,
        color="475569",
        align=PP_ALIGN.CENTER,
    )


def add_stat_hero(slide, elem: dict, resolver: TokenResolver):
    """Hero stat: very large centered number + label below."""
    x, y, w, h = _grid_rect(elem)
    accent = elem.get("accent", "#2563EB")

    val_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h * 0.6))
    vtf = val_box.text_frame
    vtf.word_wrap = True
    _add_simple_run(
        vtf,
        str(elem.get("value", "")),
        None,
        44,
        resolver,
        bold=True,
        color=accent,
        align=PP_ALIGN.CENTER,
    )

    lbl_box = slide.shapes.add_textbox(
        Inches(x), Inches(y + h * 0.6), Inches(w), Inches(h * 0.35)
    )
    ltf = lbl_box.text_frame
    ltf.word_wrap = True
    _add_simple_run(
        ltf,
        elem.get("label", ""),
        None,
        16,
        resolver,
        color="475569",
        align=PP_ALIGN.CENTER,
    )


def add_revenue_table(slide, elem: dict, resolver: TokenResolver):
    """Revenue table with name/value columns."""
    rows_data = elem.get("rows", [])
    if not rows_data:
        return
    x, y, w, h = _grid_rect(elem)
    n_rows = len(rows_data) + 1
    table_shape = slide.shapes.add_table(
        n_rows, 2, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tbl = table_shape.table
    font_name = resolver.resolve_font(None)

    for ci, header in enumerate(["Item", "Value"]):
        cell = tbl.cell(0, ci)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = font_name
                r.font.size = Pt(9)
                r.font.bold = True
                r.font.color.rgb = hex_to_rgb("1E293B")

    for ri, row in enumerate(rows_data):
        name_cell = tbl.cell(ri + 1, 0)
        name_cell.text = str(row.get("name", ""))
        val_cell = tbl.cell(ri + 1, 1)
        val_cell.text = str(row.get("value", ""))
        is_total = row.get("isTotal", False)
        for ci in range(2):
            cell = tbl.cell(ri + 1, ci)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = font_name
                    r.font.size = Pt(9)
                    r.font.bold = is_total
                    r.font.color.rgb = hex_to_rgb("1E293B" if is_total else "475569")


def add_action_list(
    slide, elem: dict, resolver: TokenResolver, slide_text_color: str | None = None
):
    """Numbered bullet list."""
    x, y, w, h = _grid_rect(elem)
    items = elem.get("items", [])
    text_color = slide_text_color or "1E293B"
    txbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        text = (
            f"{i + 1}. {item}"
            if isinstance(item, str)
            else f"{i + 1}. {item.get('text', str(item))}"
        )
        _add_simple_run(tf, text, None, 10, resolver, color=text_color)


def add_comparison(slide, elem: dict, resolver: TokenResolver):
    """Two-column comparison with headers and bullet items."""
    x, y, w, h = _grid_rect(elem)
    half = w / 2 - 0.1

    # Left column
    lhdr = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(half), Inches(0.35))
    ltf = lhdr.text_frame
    ltf.word_wrap = True
    _add_simple_run(
        ltf,
        elem.get("left_label", "Before"),
        None,
        11,
        resolver,
        bold=True,
        color="1E293B",
    )

    litems = slide.shapes.add_textbox(
        Inches(x), Inches(y + 0.4), Inches(half), Inches(h - 0.4)
    )
    lif = litems.text_frame
    lif.word_wrap = True
    for item in elem.get("left_items", []):
        if isinstance(item, dict):
            stat = item.get("stat", "")
            detail = item.get("detail", "")
            text = f"{stat}  {detail}" if stat else detail
        else:
            text = str(item)
        _add_simple_run(lif, f"• {text}", None, 9, resolver, color="475569")

    # Right column
    rx = x + half + 0.2
    rhdr = slide.shapes.add_textbox(Inches(rx), Inches(y), Inches(half), Inches(0.35))
    rtf = rhdr.text_frame
    rtf.word_wrap = True
    _add_simple_run(
        rtf,
        elem.get("right_label", "After"),
        None,
        11,
        resolver,
        bold=True,
        color="1E293B",
    )

    ritems = slide.shapes.add_textbox(
        Inches(rx), Inches(y + 0.4), Inches(half), Inches(h - 0.4)
    )
    rif = ritems.text_frame
    rif.word_wrap = True
    for item in elem.get("right_items", []):
        if isinstance(item, dict):
            stat = item.get("stat", "")
            detail = item.get("detail", "")
            text = f"{stat}  {detail}" if stat else detail
        else:
            text = str(item)
        _add_simple_run(rif, f"• {text}", None, 9, resolver, color="475569")


def add_timeline(slide, elem: dict, resolver: TokenResolver):
    """Horizontal timeline: row of date+content textboxes."""
    x, y, w, h = _grid_rect(elem)
    events = elem.get("events", [])
    if not events:
        return
    event_w = w / len(events)
    for i, evt in enumerate(events):
        ex = x + i * event_w
        box = slide.shapes.add_textbox(
            Inches(ex), Inches(y), Inches(event_w - 0.05), Inches(h)
        )
        tf = box.text_frame
        tf.word_wrap = True
        _add_simple_run(
            tf,
            str(evt.get("date", "")),
            None,
            9,
            resolver,
            bold=True,
            color="2563EB",
            align=PP_ALIGN.CENTER,
        )
        _add_simple_run(
            tf,
            str(evt.get("content", "")),
            None,
            8,
            resolver,
            color="475569",
            align=PP_ALIGN.CENTER,
        )


def add_era_card(
    slide, elem: dict, resolver: TokenResolver, slide_text_color: str | None = None
):
    """Era card: background rect + year + label + content. Dark-slide aware."""
    x, y, w, h = _grid_rect(elem)
    accent = elem.get("accent", "#2563EB")
    is_dark_slide = slide_text_color is not None

    if accent == "dim":
        card_fill = "1a3a5c"
        accent_color = "94A3B8"
    elif is_dark_slide:
        card_fill = "0d2847"
        accent_color = accent
    else:
        card_fill = "F8FAFC"
        accent_color = accent

    label_color = "F5F3EF" if is_dark_slide else "1E293B"
    body_color = "CBD5E1" if is_dark_slide else "475569"

    bg = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb(card_fill)
    bg.line.fill.background()

    ybox = slide.shapes.add_textbox(
        Inches(x + 0.1), Inches(y + 0.1), Inches(w - 0.2), Inches(0.35)
    )
    ytf = ybox.text_frame
    ytf.word_wrap = True
    _add_simple_run(
        ytf,
        str(elem.get("year", "")),
        None,
        14,
        resolver,
        bold=True,
        color=accent_color,
        align=PP_ALIGN.LEFT,
    )

    lbox = slide.shapes.add_textbox(
        Inches(x + 0.1), Inches(y + 0.45), Inches(w - 0.2), Inches(0.3)
    )
    ltf = lbox.text_frame
    ltf.word_wrap = True
    _add_simple_run(
        ltf, elem.get("label", ""), None, 10, resolver, bold=True, color=label_color
    )

    content = elem.get("content", "")
    if content:
        cbox = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + 0.8), Inches(w - 0.2), Inches(h - 0.9)
        )
        ctf = cbox.text_frame
        ctf.word_wrap = True
        _add_simple_run(ctf, content, None, 8, resolver, color=body_color)


def add_evidence_row(slide, elem: dict, resolver: TokenResolver):
    """Row of evidence textboxes."""
    x, y, w, h = _grid_rect(elem)
    items = elem.get("items", [])
    if not items:
        return
    item_w = w / len(items)
    for i, item in enumerate(items):
        ix = x + i * item_w
        box = slide.shapes.add_textbox(
            Inches(ix), Inches(y), Inches(item_w - 0.05), Inches(h)
        )
        tf = box.text_frame
        tf.word_wrap = True
        if isinstance(item, str):
            _add_simple_run(tf, item, None, 8, resolver, color="475569")
        elif isinstance(item, dict):
            label = item.get("label", "")
            detail = item.get("detail", item.get("text", item.get("content", "")))
            if label:
                _add_simple_run(
                    tf,
                    label,
                    None,
                    14,
                    resolver,
                    bold=True,
                    color="00A896",
                    align=PP_ALIGN.CENTER,
                )
            if detail:
                _add_simple_run(
                    tf,
                    str(detail),
                    None,
                    8,
                    resolver,
                    color="475569",
                    align=PP_ALIGN.CENTER,
                )
        else:
            _add_simple_run(tf, str(item), None, 8, resolver, color="475569")


def add_evidence_proof(slide, elem: dict, resolver: TokenResolver):
    """2-column evidence image grid with captions."""
    x, y, w, h = _grid_rect(elem)
    images = elem.get("images", [])
    if not images:
        return
    col_w = (w - 0.2) / 2
    for i, img in enumerate(images[:2]):
        ix = x + i * (col_w + 0.2)
        src = img.get("src", "")
        if Path(src).exists():
            slide.shapes.add_picture(
                str(src), Inches(ix), Inches(y), Inches(col_w), Inches(h * 0.75)
            )
        caption = img.get("caption", "")
        if caption:
            cbox = slide.shapes.add_textbox(
                Inches(ix), Inches(y + h * 0.78), Inches(col_w), Inches(h * 0.2)
            )
            ctf = cbox.text_frame
            ctf.word_wrap = True
            _add_simple_run(ctf, caption, None, 7, resolver, color="6B7280")


def add_radial(slide, elem: dict, resolver: TokenResolver):
    """Hub-and-spoke radial layout degraded to positioned shapes."""
    from pptx.enum.shapes import MSO_SHAPE

    x, y, w, h = _grid_rect(elem)
    hub = elem.get("hub", {})
    spokes = elem.get("spokes", [])

    cx, cy = x + w / 2, y + h / 2
    hub_r = 0.8
    hub_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - hub_r),
        Inches(cy - hub_r),
        Inches(hub_r * 2),
        Inches(hub_r * 2),
    )
    hub_shape.fill.solid()
    hub_shape.fill.fore_color.rgb = resolver.resolve_color("00356B")
    hub_shape.line.fill.background()
    tf = hub_shape.text_frame
    tf.word_wrap = True
    _add_simple_run(
        tf,
        hub.get("label", ""),
        None,
        9,
        resolver,
        bold=True,
        color="FFFFFF",
        align=PP_ALIGN.CENTER,
    )

    positions = {
        "top": (cx - 1.0, cy - hub_r - 1.2),
        "right": (cx + hub_r + 0.4, cy - 0.3),
        "bottom-right": (cx + hub_r + 0.1, cy + hub_r + 0.3),
        "bottom-left": (cx - hub_r - 2.0, cy + hub_r + 0.3),
        "left": (cx - hub_r - 2.3, cy - 0.3),
    }

    for spoke in spokes:
        pos = spoke.get("position", "right")
        sx, sy = positions.get(pos, (cx + 1.5, cy))
        node_w, node_h = 2.0, 0.6
        node = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(sx),
            Inches(sy),
            Inches(node_w),
            Inches(node_h),
        )
        node.fill.solid()
        node.fill.fore_color.rgb = resolver.resolve_color("FFFFFF")
        node.line.color.rgb = resolver.resolve_color("E5E2DC")
        ntf = node.text_frame
        ntf.word_wrap = True
        label = spoke.get("label", "")
        stat = spoke.get("stat", "")
        _add_simple_run(
            ntf,
            f"{label}\n{stat}",
            None,
            8,
            resolver,
            bold=True,
            color="00356B",
            align=PP_ALIGN.CENTER,
        )


# Component Type (HTML)    → Handler (PPTX)        → Notes
# title/content            → add_textbox
# two-column               → add_comparison
# data/data-annotation     → add_image (chart SVG)  / add_shape
# stat/hero-stat           → add_stat_card/hero
# quote                    → add_textbox
# full-bleed-image         → add_image
# timeline                 → add_timeline
# before-after             → add_comparison
# evidence-cascade         → add_evidence_row
# section-break            → add_textbox
# storyboard               → add_evidence_proof
# argument-flow            → add_action_list
# overlapping-cascade      → add_evidence_row       (reuse, no rotation in PPTX)
ELEMENT_HANDLERS = {
    "text": add_textbox,
    "shape": add_shape,
    "image": add_image,
    "table": add_table,
    "card": add_card,
    "stat_card": add_stat_card,
    "stat_hero": add_stat_hero,
    "revenue_table": add_revenue_table,
    "action_list": add_action_list,
    "comparison": add_comparison,
    "timeline": add_timeline,
    "era_card": add_era_card,
    "evidence_row": add_evidence_row,
    "evidence_proof": add_evidence_proof,
    "radial": add_radial,
    "overlapping_cascade": add_evidence_row,
}


TITLE_HEIGHT = 0.7
CONTENT_START_Y = 1.0
ELEMENT_GAP = 0.15

DEFAULT_HEIGHTS: dict[str, float] = {
    "text": 0.5,
    "card": 1.8,
    "stat_card": 1.4,
    "stat_hero": 1.8,
    "table": 2.5,
    "revenue_table": 2.2,
    "comparison": 2.5,
    "timeline": 0.8,
    "era_card": 1.8,
    "evidence_row": 0.7,
    "evidence_proof": 1.6,
    "action_list": 1.5,
    "radial": 3.0,
    "image": 2.0,
    "shape": 1.0,
}


HANDLERS_WITH_SLIDE_TEXT_COLOR = {"text", "action_list", "era_card"}

TITLE_ROLE_POSITIONS = {
    "hero_title": (SLIDE_HEIGHT * 0.30, 0.9, 36, True),
    "hero_subtitle": (SLIDE_HEIGHT * 0.52, 0.5, 18, False),
    "meta": (SLIDE_HEIGHT * 0.68, 0.35, 11, False),
    "title": (SLIDE_HEIGHT * 0.25, 0.8, 32, True),
    "cta": (SLIDE_HEIGHT * 0.72, 0.4, 16, True),
    "close_quote": (SLIDE_HEIGHT * 0.78, 0.8, 10, False),
}


def _estimate_min_height(
    content: str, font_size_pt: float, box_w_inches: float
) -> float:
    """Estimate minimum height needed for text to render without clipping."""
    if not content:
        return 0.5
    avg_char_width_in = font_size_pt / 72 * 0.6
    chars_per_line = max(1, box_w_inches / avg_char_width_in)
    _math = __import__("math")
    num_lines = _math.ceil(len(content) / chars_per_line)
    line_height_in = font_size_pt / 72 * 1.3
    return max(0.5, num_lines * line_height_in + 0.1)


def _check_collisions(
    placed: list[tuple[str, float, float, float, float]],
) -> list[str]:
    """Check for overlapping placed elements. Returns warning strings."""
    warnings = []
    for i, (desc_a, x_a, y_a, w_a, h_a) in enumerate(placed):
        for j, (desc_b, x_b, y_b, w_b, h_b) in enumerate(placed[i + 1 :], i + 1):
            ox = max(0, min(x_a + w_a, x_b + w_b) - max(x_a, x_b))
            oy = max(0, min(y_a + h_a, y_b + h_b) - max(y_a, y_b))
            overlap = ox * oy
            smaller = min(w_a * h_a, w_b * h_b)
            if smaller > 0 and overlap / smaller > 0.10:
                warnings.append(
                    f"Collision: {desc_a} overlaps {desc_b} ({overlap / smaller:.0%})"
                )
    return warnings


def build_slide(
    prs: PresentationClass, slide_data: dict, resolver: TokenResolver
) -> dict:
    """Build a single slide from structured data. Returns element count for validation.

    Maintains a running y_cursor so elements stack vertically instead of
    overlapping at a hardcoded y=1.2. Elements with explicit 'y' values
    in the JSON override the cursor; elements without 'y' use the cursor
    and advance it by their height + gap.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide_data.get("background")
    set_background(slide, bg, resolver)

    slide_text_color = slide_data.get("textColor")
    slide_type = slide_data.get("type", "content")
    is_title_slide = slide_type in ("title", "close")

    elements = slide_data.get("elements", [])
    built = 0
    y_cursor = CONTENT_START_Y
    row_max_bottom = y_cursor
    row_cols_used: set[int] = set()
    placed_elements: list[tuple[str, float, float, float, float]] = []

    # Pre-scan: find max bottom of all role-positioned elements on title/close slides
    role_max_bottom = 0.0
    if is_title_slide:
        for el in elements:
            if el.get("type", "text") == "text":
                rp = TITLE_ROLE_POSITIONS.get(el.get("role", ""))
                if rp:
                    role_max_bottom = max(role_max_bottom, rp[0] + rp[1])

    for elem in elements:
        elem_type = elem.get("type", "text")
        handler = ELEMENT_HANDLERS.get(elem_type)

        if handler:
            has_explicit_y = "y" in elem
            has_explicit_h = "h" in elem

            # Title/close slides: center text elements vertically
            if is_title_slide and elem_type == "text":
                role = elem.get("role", "")
                pos = TITLE_ROLE_POSITIONS.get(role)
                if pos:
                    ty, th, fsize, bold = pos
                    elem = {
                        **elem,
                        "x": MARGIN,
                        "y": ty,
                        "w": SLIDE_WIDTH - 2 * MARGIN,
                        "h": th,
                        "font_size": fsize,
                        "bold": bold,
                        "align": "center",
                    }
                    if slide_text_color and "color" not in elem:
                        elem = {**elem, "color": slide_text_color}
                    has_explicit_y = True
                    has_explicit_h = True
                    # Advance cursor past this positioned element
                    pos_bottom = ty + th + ELEMENT_GAP
                    if pos_bottom > y_cursor:
                        y_cursor = pos_bottom
                        row_max_bottom = y_cursor

            # S6 fix: nudge originally-explicit-y elements below role positions
            orig_had_y = "y" in elem and not (
                is_title_slide
                and elem_type == "text"
                and TITLE_ROLE_POSITIONS.get(elem.get("role", ""))
            )
            if orig_had_y and is_title_slide and role_max_bottom > 0:
                if CONTENT_START_Y < elem["y"] < role_max_bottom + 0.1:
                    new_y = role_max_bottom + 0.15
                    log.warning(
                        "RECOVER: nudging element from y=%.2f to y=%.2f"
                        " (below role positions)",
                        elem["y"],
                        new_y,
                    )
                    elem = {**elem, "y": new_y}

            if not has_explicit_h:
                default_h = DEFAULT_HEIGHTS.get(elem_type, 0.5)
                # Dynamic height for tables based on row count
                if elem_type == "table":
                    headers = elem.get("headers", [])
                    rows_data = elem.get("rows", [])
                    n_rows = len(rows_data) + (1 if headers else 0)
                    has_multiline = any(
                        "\n" in str(cell)
                        for row in rows_data
                        for cell in (row if isinstance(row, list) else [row])
                    )
                    row_h = 0.55 if has_multiline else 0.4
                    default_h = max(2.5, n_rows * row_h)
                # Dynamic height for text based on content length
                if elem_type == "text" and default_h == 0.5:
                    content = elem.get("content", "")
                    if len(content) > 200:
                        default_h = 1.0
                    elif len(content) > 100:
                        default_h = 0.7
                # Title roles: content-aware height using font size + box width
                if elem_type == "text" and elem.get("role") in (
                    "title",
                    "hero_title",
                ):
                    content = elem.get("content", "")
                    font_pt = elem.get("font_size", 20)
                    box_w = elem.get("w", SLIDE_WIDTH - 2 * MARGIN)
                    min_h = _estimate_min_height(content, font_pt, box_w)
                    default_h = max(default_h, min_h)
                elem = {**elem, "h": default_h}

            if not has_explicit_y:
                gc = elem.get("gridColumn", "1 / -1")
                parts = [p.strip() for p in gc.split("/")]
                col_start = int(parts[0]) - 1
                col_end = 12 if parts[1].strip() == "-1" else int(parts[1]) - 1
                elem_cols = set(range(col_start, col_end))

                if elem_cols & row_cols_used:
                    y_cursor = row_max_bottom + ELEMENT_GAP
                    row_cols_used = set()
                    row_max_bottom = y_cursor

                elem = {**elem, "y": y_cursor}
                row_cols_used |= elem_cols
                elem_bottom = y_cursor + elem["h"]
                if elem_bottom > row_max_bottom:
                    row_max_bottom = elem_bottom

            # Off-slide overflow guard
            ey = elem.get("y", y_cursor)
            eh = elem.get("h", 0.5)
            if ey + eh > SLIDE_HEIGHT - 0.1:
                log.warning(
                    "RECOVER: element '%s' bottom %.2f\" exceeds slide height %.2f\"",
                    str(elem.get("content", elem_type))[:40],
                    ey + eh,
                    SLIDE_HEIGHT,
                )

            # Track for collision detection — resolve grid position
            desc = f"{elem_type}:{str(elem.get('content', elem.get('value', '')))[:25]}"
            gx, _, gw, _ = _grid_rect(elem)
            ex = elem.get("x", gx)
            ew = elem.get("w", gw)
            placed_elements.append((desc, ex, ey, ew, eh))

            if elem_type == "image":
                handler(slide, elem)
            elif elem_type in HANDLERS_WITH_SLIDE_TEXT_COLOR:
                handler(slide, elem, resolver, slide_text_color=slide_text_color)
            else:
                handler(slide, elem, resolver)
            built += 1
        else:
            log.warning("Unknown element type: %s", elem_type)

    # Post-layout collision detection
    for warning in _check_collisions(placed_elements):
        log.warning("LAYOUT: %s (slide type=%s)", warning, slide_type)

    return {
        "type": slide_data.get("type", "unknown"),
        "elements_built": built,
        "elements_total": len(elements),
    }


def validate_roundtrip(
    pptx_path: str, expected_slides: int, expected_elements: list[int]
):
    """Re-open PPTX and verify slide/element counts. E-PPTX-002 safety check."""
    try:
        prs = Presentation(pptx_path)
        actual_slides = len(prs.slides)

        if actual_slides != expected_slides:
            log.error(
                "Round-trip FAIL: expected %d slides, got %d",
                expected_slides,
                actual_slides,
            )
            return False

        for i, slide in enumerate(prs.slides):
            actual_shapes = len(slide.shapes)
            if i < len(expected_elements) and actual_shapes < expected_elements[i]:
                log.warning(
                    "Slide %d: expected ≥%d shapes, got %d",
                    i + 1,
                    expected_elements[i],
                    actual_shapes,
                )

        log.info("Round-trip OK: %d slides verified", actual_slides)
        return True
    except Exception as e:
        log.error("Round-trip validation failed: %s", e)
        return False


def build_pptx(
    slides_json: dict, tokens: dict | None = None, output_path: str | None = None
) -> str:
    """Build PPTX from sidecar JSON + optional design tokens."""
    resolver = TokenResolver(tokens)

    metadata = slides_json.get("metadata", {})
    slides = slides_json.get("slides", [])

    if not slides:
        log.error("No slides found in sidecar JSON")
        sys.exit(1)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)

    inline_tokens = slides_json.get("design_tokens")
    if inline_tokens and not tokens:
        resolver = TokenResolver(inline_tokens)

    slide_stats = []
    for slide_data in slides:
        stats = build_slide(prs, slide_data, resolver)
        slide_stats.append(stats)

    if not output_path:
        slug = metadata.get("title", "presentation").lower().replace(" ", "-")[:40]
        output_path = f"output/pptx/{slug}-{date.today()}.pptx"

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))

    expected_elements = [s["elements_built"] for s in slide_stats]
    validate_roundtrip(str(out), len(slides), expected_elements)

    font_warnings = resolver.fallback_map
    if font_warnings:
        log.warning("Font fallbacks applied: %s", font_warnings)

    return str(out)


def build_pptx_via_node(json_path: str, output_path: str | None = None) -> str:
    """Engine B: delegate to PptxGenJS via Node.js subprocess."""
    import subprocess  # noqa: PLC0415 — lazy import, formatter strips top-level

    engine_dir = Path(__file__).parent / "pptxgen"
    engine_script = engine_dir / "engine.mjs"

    if not engine_script.exists():
        log.error("PptxGenJS engine not found at %s", engine_script)
        sys.exit(1)

    cmd = ["node", str(engine_script), json_path]
    if output_path:
        cmd.append(output_path)

    result = subprocess.run(cmd, capture_output=True, text=True, cwd=str(engine_dir))

    if result.returncode != 0:
        log.error("PptxGenJS engine failed:\n%s", result.stderr)
        sys.exit(1)

    for line in result.stdout.splitlines():
        if line.startswith("RESULT_JSON:"):
            data = json.loads(line[len("RESULT_JSON:") :])
            return data["output"]

    log.error("PptxGenJS engine produced no RESULT_JSON line")
    sys.exit(1)


def main():
    logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")

    if len(sys.argv) < 2:
        print(
            "Usage: slides_to_pptx.py <slides.json> [--tokens <ref.yaml>] [--out <output.pptx>] [--engine python-pptx|pptxgenjs]",
            file=sys.stderr,
        )
        sys.exit(1)

    json_path = sys.argv[1]
    tokens_path = None
    output_path = None
    engine = "python-pptx"

    args = sys.argv[2:]
    i = 0
    while i < len(args):
        if args[i] == "--tokens" and i + 1 < len(args):
            tokens_path = args[i + 1]
            i += 2
        elif args[i] == "--out" and i + 1 < len(args):
            output_path = args[i + 1]
            i += 2
        elif args[i] == "--engine" and i + 1 < len(args):
            engine = args[i + 1]
            i += 2
        else:
            i += 1

    if engine == "pptxgenjs":
        out = build_pptx_via_node(json_path, output_path)
        print(f"PPTX saved: {out}")
        with open(json_path) as f:
            print(f"Slides: {len(json.load(f).get('slides', []))}")
        return

    if engine != "python-pptx":
        print(
            f"Unknown engine '{engine}'. Use --engine python-pptx or --engine pptxgenjs",
            file=sys.stderr,
        )
        sys.exit(1)

    with open(json_path) as f:
        slides_json = json.load(f)

    tokens = None
    if tokens_path:
        with open(tokens_path) as f:
            tokens = yaml.safe_load(f)

    out = build_pptx(slides_json, tokens, output_path)
    print(f"PPTX saved: {out}")
    print(f"Slides: {len(slides_json.get('slides', []))}")


if __name__ == "__main__":
    main()
