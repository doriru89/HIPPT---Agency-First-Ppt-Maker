#!/usr/bin/env python3
"""Extract layout library from existing PPTX files.

Converts PPTX shapes to percentage-based layout regions using positional,
font-based, and placeholder-type heuristics (no CSS selectors available).
Output is identical in schema to extract_layouts.py (HTML path).

Usage:
    uv run python TOOLS/scripts/pptx_to_layout.py file.pptx
    uv run python TOOLS/scripts/pptx_to_layout.py file.pptx --out layouts/ --verify
"""

import argparse
import logging
import re
from collections import defaultdict
from datetime import date
from pathlib import Path


log = logging.getLogger(__name__)

from hippt.layout_utils import (
    auto_supports,
    auto_tags,
    cluster_by_x,
    dominant_role,
    element_type,
    get_thresholds,
    hex_color,
    infer_slide_type,
    load_config,
    merge_bbox,
    pct_box,
    split_y_bands,
    verify_layouts,
    write_layout_index,
    write_layout_yaml,
)

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

from hippt.design_tokens import VIEWPORT_H, VIEWPORT_W

# ── PPTX element adapter ────────────────────────────────────────────────────


def _extract_text_props(shape):
    """Extract (text, max_font_pt, is_bold, font_name) from a shape."""
    if not shape.has_text_frame:
        return "", 0.0, False, ""
    text_parts = []
    max_fs = 0.0
    bold = False
    name = ""
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            text_parts.append(run.text)
            if run.font.size is not None:
                fs = run.font.size / Pt(1)
                if fs > max_fs:
                    max_fs = fs
                    bold = bool(run.font.bold)
                    name = run.font.name or ""
    return "".join(text_parts).strip(), max_fs, bold, name


def _get_fill_hex(shape):
    """Get shape fill color as hex, or None if transparent/inherited."""
    try:
        fill = shape.fill
        if fill and fill.type is not None:
            return hex_color(fill.fore_color.rgb)
    except Exception as exc:
        log.debug("Skipping fill extraction: %s", exc)
    return None


def _shape_to_element(shape, slide_w_emu, slide_h_emu, vw, vh):
    """Convert a python-pptx Shape to the normalized element dict format."""
    left = shape.left or 0
    top = shape.top or 0
    width = shape.width or 0
    height = shape.height or 0

    if not slide_w_emu or not slide_h_emu:
        return None

    x_px = left / slide_w_emu * vw
    y_px = top / slide_h_emu * vh
    w_px = width / slide_w_emu * vw
    h_px = height / slide_h_emu * vh

    bbox = {"x": x_px, "y": y_px, "w": w_px, "h": h_px}

    st = None
    try:
        st = shape.shape_type
    except Exception as exc:
        log.debug("Cannot read shape_type: %s", exc)

    if st == MSO_SHAPE_TYPE.PICTURE:
        etype = "image"
    elif st == MSO_SHAPE_TYPE.TABLE:
        etype = "table"
    elif st == MSO_SHAPE_TYPE.CHART:
        etype = "table"
    elif st in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
        etype = "shape"
    else:
        etype = "text"

    text, font_pt, bold, font_name = _extract_text_props(shape)
    is_decorative = (h_px < 3 or w_px < 3) and not text

    ph_idx = None
    try:
        if shape.placeholder_format:
            ph_idx = shape.placeholder_format.idx
    except (AttributeError, ValueError):
        pass

    return {
        "type": etype,
        "bbox": bbox,
        "text": text,
        "style": {
            "fontSize": f"{font_pt}px" if font_pt else "0px",
            "fontWeight": "700" if bold else "400",
            "fontFamily": font_name,
        },
        "selector": "",
        "isDecorative": is_decorative,
        "src": None,
        "_ph_idx": ph_idx,
        "_font_pt": font_pt,
        "_fill_color": _get_fill_hex(shape),
    }


def _flatten_group(group_shape, slide_w_emu, slide_h_emu, vw, vh):
    """Recursively expand a GroupShape into a flat element list."""
    elements = []
    for child in group_shape.shapes:
        try:
            if child.shape_type == MSO_SHAPE_TYPE.GROUP:
                elements.extend(_flatten_group(child, slide_w_emu, slide_h_emu, vw, vh))
            else:
                elem = _shape_to_element(child, slide_w_emu, slide_h_emu, vw, vh)
                if elem:
                    elements.append(elem)
        except Exception as exc:
            log.debug("Skipping grouped shape: %s", exc)
    return elements


# ── Role classification ─────────────────────────────────────────────────────


def _classify_role_pptx(elem, cfg):
    """7-level role cascade for PPTX elements (no CSS selectors)."""
    pptx_cfg = cfg.get("pptx_roles", {})
    vw, vh = VIEWPORT_W, VIEWPORT_H

    # Level 1: Placeholder index
    ph_idx = elem.get("_ph_idx")
    if ph_idx is not None:
        ph_map = pptx_cfg.get("placeholder_map", {})
        role = ph_map.get(ph_idx)
        if role:
            return role

    # Level 2: Shape type
    etype = elem.get("type")
    if etype == "image":
        return "evidence_image"
    if etype == "table":
        return "data_table"

    # Level 3: Decorative detection
    bbox = elem["bbox"]
    text = (elem.get("text") or "").strip()
    if elem.get("isDecorative"):
        return "divider" if bbox["h"] < 6 else "decorative_shape"
    area_pct = (bbox["w"] * bbox["h"]) / (vw * vh) * 100
    if area_pct < 2 and not text:
        return "decorative_shape"
    if bbox["h"] < 6 and bbox["w"] > vw * 0.3:
        return "divider"

    # Level 4: Positional heuristics
    pos = pptx_cfg.get("position_rules", {})
    y_pct = bbox["y"] / vh * 100
    w_pct = bbox["w"] / vw * 100

    if y_pct < pos.get("progress_bar_y_max_pct", 2) and w_pct > pos.get(
        "progress_bar_min_width_pct", 50
    ):
        if bbox["h"] < 10:
            return "progress_bar"

    if y_pct > pos.get("footer_y_min_pct", 88) and text and len(text) < 60:
        return "footer"

    if y_pct < pos.get("section_label_y_max_pct", 5) and text and len(text) < 30:
        font_pt = elem.get("_font_pt", 0)
        if 0 < font_pt < 14:
            return "section_label"

    # Level 5: Text content patterns
    if text:
        for p in cfg.get("text_patterns", []):
            if re.match(p["pattern"], text):
                return p["role"]
        for p in pptx_cfg.get("text_patterns_pptx", []):
            conds = p.get("conditions", {})
            if re.match(p["pattern"], text):
                min_fp = conds.get("min_font_pt", 0)
                if elem.get("_font_pt", 0) >= min_fp:
                    return p["role"]

    # Level 6: Font-size classification
    font_pt = elem.get("_font_pt", 0)
    ft = pptx_cfg.get("font_thresholds", {})
    if font_pt > ft.get("headline_min_pt", 24):
        return "headline"
    sh_range = ft.get("section_heading_range", [15, 17])
    fw = elem.get("style", {}).get("fontWeight", "400")
    if sh_range[0] <= font_pt <= sh_range[1] and fw in ("700", "800", "bold"):
        return "section_heading"
    if font_pt <= ft.get("detail_text_max_pt", 12) and font_pt > 0:
        return "detail_text"

    # Level 7: Default
    return "body_text"


# ── Grouping ────────────────────────────────────────────────────────────────


def _style_family_key(elem, bucket_pt):
    """Build a grouping key from font size bucket + weight + fill color."""
    fs = elem.get("_font_pt", 0)
    bucketed = round(fs / bucket_pt) * bucket_pt if bucket_pt else fs
    fw = elem.get("style", {}).get("fontWeight", "400")
    fill = elem.get("_fill_color") or "none"
    return (bucketed, fw, fill)


def _group_elements_pptx(elements, vw, vh, cfg):
    """Group PPTX elements using style-family bucketing instead of CSS families."""
    th = get_thresholds(cfg)
    structural_roles = set(
        cfg.get("structural_roles", ["progress_bar", "section_label", "slide_counter"])
    )
    merge_roles = cfg.get("merge_roles", {})
    pptx_cfg = cfg.get("pptx_roles", {})
    sg = pptx_cfg.get("style_grouping", {})
    bucket_pt = sg.get("font_size_bucket_pt", 2)
    isolated = set(sg.get("isolated_roles", ["headline", "subtitle", "radial_hub"]))

    x_thresh = vw * th["x_prox"]
    y_band_thresh = vh * th["y_band"]
    generic_max = vh * th["generic_max"]
    x_margin = vw * th["x_margin"]

    regions = []

    # Phase 1: structural isolation
    progress = [e for e in elements if e["_role"] == "progress_bar"]
    if progress:
        regions.append(
            {
                "role": "progress_bar",
                "bbox": merge_bbox([e["bbox"] for e in progress]),
                "elements": progress,
                "element_type": "shape",
            }
        )
    for e in elements:
        if e["_role"] in structural_roles and e["_role"] != "progress_bar":
            regions.append(
                {
                    "role": e["_role"],
                    "bbox": e["bbox"],
                    "elements": [e],
                    "element_type": "text",
                }
            )

    content = [e for e in elements if e["_role"] not in structural_roles]

    # Phase 2: style-family grouping
    family_groups = defaultdict(list)
    isolated_elems = []
    for e in content:
        if e["_role"] in isolated:
            isolated_elems.append(e)
        else:
            key = _style_family_key(e, bucket_pt)
            family_groups[key].append(e)

    content_regions = []

    for e in isolated_elems:
        et = "image" if e.get("type") == "image" else "text"
        content_regions.append(
            {
                "role": e["_role"],
                "bbox": e["bbox"],
                "elements": [e],
                "element_type": et,
            }
        )

    # Phase 3: X-clustering within families
    for _key, elems in family_groups.items():
        for xc in cluster_by_x(elems, x_thresh):
            content_regions.append(
                {
                    "role": dominant_role(xc, cfg),
                    "bbox": merge_bbox([e["bbox"] for e in xc]),
                    "elements": list(xc),
                    "element_type": element_type(xc),
                }
            )

    # Phase 4: generic proximity pass (single-element regions try to attach)
    settled = [r for r in content_regions if len(r["elements"]) > 1]
    singles = [r for r in content_regions if len(r["elements"]) == 1]
    for sr in sorted(singles, key=lambda r: r["bbox"]["y"]):
        g = sr["elements"][0]
        gx = g["bbox"]["x"] + g["bbox"]["w"] / 2
        gy = g["bbox"]["y"] + g["bbox"]["h"] / 2
        best_cr, best_dist = None, float("inf")
        for cr in settled:
            cb = cr["bbox"]
            rx1 = cb["x"] - x_margin
            rx2 = cb["x"] + cb["w"] + x_margin
            if not (rx1 <= gx <= rx2):
                continue
            ry_bot = cb["y"] + cb["h"]
            dist = max(0, gy - ry_bot) if gy > ry_bot else max(0, cb["y"] - gy)
            if dist < best_dist:
                best_dist, best_cr = dist, cr
        if best_cr is not None and best_dist < generic_max:
            best_cr["elements"].append(g)
            best_cr["bbox"] = merge_bbox([e["bbox"] for e in best_cr["elements"]])
            best_cr["element_type"] = element_type(best_cr["elements"])
        else:
            settled.append(sr)
    content_regions = settled

    # Phase 5: Y-band splitting
    min_split = th["min_split"]
    split_regions = []
    for cr in content_regions:
        if len(cr["elements"]) >= min_split:
            bands = split_y_bands(cr["elements"], y_band_thresh)
        else:
            bands = [cr["elements"]]
        for band in bands:
            split_regions.append(
                {
                    "role": dominant_role(band, cfg),
                    "bbox": merge_bbox([e["bbox"] for e in band]),
                    "elements": band,
                    "element_type": element_type(band),
                }
            )
    content_regions = split_regions

    # Phase 6: post-merge
    for src_role, merged_name in merge_roles.items():
        matches = [r for r in content_regions if r["role"] == src_role]
        if len(matches) > 1:
            all_elems = [e for r in matches for e in r["elements"]]
            kept = [r for r in content_regions if r["role"] != src_role]
            kept.append(
                {
                    "role": merged_name,
                    "bbox": merge_bbox([e["bbox"] for e in all_elems]),
                    "elements": all_elems,
                    "element_type": element_type(all_elems),
                }
            )
            content_regions = kept

    regions.extend(content_regions)
    regions.sort(key=lambda r: (r["bbox"]["y"], r["bbox"]["x"]))
    return regions


# ── Slide-level extraction ──────────────────────────────────────────────────


def _extract_slide_bg(slide):
    """Determine slide background as 'light' or 'dark'."""
    try:
        bg = slide.background
        if bg and bg.fill and bg.fill.type is not None:
            rgb = bg.fill.fore_color.rgb
            r, g, b = rgb[0], rgb[1], rgb[2]
            lum = (0.299 * r + 0.587 * g + 0.114 * b) / 255
            return "dark" if lum < 0.5 else "light"
    except Exception as exc:
        log.debug("Slide bg extraction failed: %s", exc)
    try:
        layout_bg = slide.slide_layout.background
        if layout_bg and layout_bg.fill and layout_bg.fill.type is not None:
            rgb = layout_bg.fill.fore_color.rgb
            r, g, b = rgb[0], rgb[1], rgb[2]
            lum = (0.299 * r + 0.587 * g + 0.114 * b) / 255
            return "dark" if lum < 0.5 else "light"
    except Exception as exc:
        log.debug("Layout bg extraction failed: %s", exc)
    return "light"


_TYPE_PREFIX = {
    "title": "TITLE",
    "close": "CLOSE",
    "data": "DATA",
    "editorial": "EDIT",
    "comparison": "COMP",
    "radial": "RADIAL",
    "timeline": "ERA",
    "hero": "HERO",
}


def extract_pptx_slide_layout(slide, index, prs, source="pptx", config_path=None):
    """Extract one layout dict from a python-pptx Slide."""
    cfg = load_config(config_path)
    vw, vh = VIEWPORT_W, VIEWPORT_H
    slide_w_emu = prs.slide_width
    slide_h_emu = prs.slide_height

    elements = []
    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                elements.extend(_flatten_group(shape, slide_w_emu, slide_h_emu, vw, vh))
            else:
                elem = _shape_to_element(shape, slide_w_emu, slide_h_emu, vw, vh)
                if elem:
                    elements.append(elem)
        except Exception as exc:
            log.debug("Skipping shape %s: %s", getattr(shape, "name", "?"), exc)

    if not elements:
        return None

    for e in elements:
        e["_role"] = _classify_role_pptx(e, cfg)

    raw_regions = _group_elements_pptx(elements, vw, vh, cfg)
    structural = set(
        cfg.get("structural_roles", ["progress_bar", "section_label", "slide_counter"])
    )
    content_count = sum(1 for r in raw_regions if r["role"] not in structural)

    regions_out = []
    for r in raw_regions:
        regions_out.append(
            {
                "role": r["role"],
                "element_type": r["element_type"],
                "box": pct_box(r["bbox"], vw, vh),
            }
        )

    total_slides = len(prs.slides)
    master_name = ""
    try:
        master_name = slide.slide_layout.name
    except Exception as exc:
        log.debug("Cannot read slide layout name: %s", exc)

    slide_type = infer_slide_type(raw_regions, index, total_slides, master_name)
    prefix = _TYPE_PREFIX.get(slide_type, "UNK")
    code = f"L-{prefix}-{index + 1:03d}"
    name = f"{slide_type.title()} (PPTX s{index + 1})"

    background = _extract_slide_bg(slide)
    density = (
        "low" if content_count <= 3 else ("medium" if content_count <= 6 else "high")
    )

    ar_w = slide_w_emu / slide_h_emu
    aspect = "16:9" if abs(ar_w - 16 / 9) < 0.05 else f"{ar_w:.2f}:1"

    return {
        "code": code,
        "name": name,
        "source": f"{source}-s{index + 1}",
        "slide_type": slide_type,
        "background": background,
        "density": density,
        "max_regions": content_count,
        "tags": auto_tags(raw_regions),
        "supports": auto_supports(raw_regions),
        "regions": regions_out,
        "origin": "pptx",
        "slide_master": master_name or None,
        "aspect_ratio": aspect,
        "thumbnail": None,
        "context": "",
        "anti_patterns": "",
        "created": str(date.today()),
        "verified_in": "pending",
        "calibration_status": "candidate",
    }


# ── File-level API ──────────────────────────────────────────────────────────


def extract_pptx_layouts(
    pptx_path,
    source=None,
    output_dir=None,
    config_path=None,
    verify=False,
):
    """Extract layout YAMLs from a PPTX file.

    Returns list of layout dicts (same schema as extract_layouts.py output).
    If output_dir is given, writes L-*.yaml + index.yaml there.
    """
    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))
    src = source or pptx_path.stem

    layouts = []
    for i, slide in enumerate(prs.slides):
        layout = extract_pptx_slide_layout(slide, i, prs, src, config_path)
        if layout:
            layouts.append(layout)
            log.info(
                "S%d → %s (%d regions)", i + 1, layout["code"], len(layout["regions"])
            )

    if output_dir:
        out = Path(output_dir)
        out.mkdir(parents=True, exist_ok=True)
        for la in layouts:
            write_layout_yaml(la, out)
        write_layout_index(layouts, out)
        log.info("Index: %d layouts → %s/index.yaml", len(layouts), out)

    if verify:
        viewport = {"w": VIEWPORT_W, "h": VIEWPORT_H}
        warns = verify_layouts(layouts, viewport)
        for w in warns:
            log.warning("VERIFY: %s", w)
        if not warns:
            log.info("VERIFY: all regions within tolerance")

    return layouts


# ── CLI ─────────────────────────────────────────────────────────────────────


def main():
    parser = argparse.ArgumentParser(
        description="Extract layout library from PPTX files"
    )
    parser.add_argument("pptx_path", help="Path to input PPTX file")
    parser.add_argument("--out", default="output/layouts")
    parser.add_argument("--source", default=None)
    parser.add_argument("--verify", action="store_true")
    parser.add_argument("--config", default=None)
    args = parser.parse_args()

    logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")

    extract_pptx_layouts(
        args.pptx_path,
        source=args.source,
        output_dir=args.out,
        config_path=args.config,
        verify=args.verify,
    )


if __name__ == "__main__":
    main()
