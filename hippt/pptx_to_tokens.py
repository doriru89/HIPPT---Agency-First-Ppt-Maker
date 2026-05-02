#!/usr/bin/env python3
"""Convert PPTX analysis to design tokens YAML matching reference-extract schema.

Reads a PPTX file directly, runs analyze_pptx for metrics, adds font family
extraction, classifies colors/typography by role, computes spacing rhythm
and hierarchy weight ratio. Outputs YAML compatible with /reference-extract.

Usage:
    uv run python TOOLS/scripts/pptx_to_tokens.py <file.pptx> [--slug NAME]
"""

from __future__ import annotations

import math
import sys
from collections import Counter
from datetime import date
from pathlib import Path

import yaml
from pptx import Presentation

from hippt.analyze_pptx import analyze_file


def extract_font_families(pptx_path: str) -> list[dict]:
    """Extract font family names with size and frequency from PPTX runs."""
    prs = Presentation(pptx_path)
    fonts: Counter[tuple[str, float]] = Counter()

    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        name = run.font.name
                        size = run.font.size
                        if name and size:
                            from pptx.util import Pt

                            fonts[(name, round(size / Pt(1), 1))] += 1
            except Exception:
                continue

    result = []
    for (family, size), count in fonts.most_common():
        result.append({"family": family, "size": f"{size}pt", "frequency": count})
    return result


def classify_color_role(hex_val: str, rank: int, total_colors: int, count: int) -> str:
    """Classify color role by frequency rank and luminance."""
    r = int(hex_val[1:3], 16)
    g = int(hex_val[3:5], 16)
    b = int(hex_val[5:7], 16)
    luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255

    freq_pct = count / total_colors if total_colors > 0 else 0

    if rank == 0 and luminance > 0.7:
        return "bg"
    if rank == 0 and luminance < 0.3:
        return "text"
    if freq_pct > 0.3:
        return "bg" if luminance > 0.5 else "text"
    if freq_pct < 0.05:
        saturation = (max(r, g, b) - min(r, g, b)) / max(max(r, g, b), 1)
        if saturation > 0.3:
            return "accent"
        if luminance > 0.85:
            return "bg"
        return "surface"
    if luminance > 0.5:
        return "surface"
    return "text"


from hippt.design_tokens import classify_font_role  # noqa: E402, PLC0415


def compute_spacing_rhythm(padding_values: list[float]) -> tuple[str, list[int]]:
    """Find spacing rhythm from padding percentages converted to approximate px."""
    # Convert padding % to px using canonical viewport width (E-PRES-036)
    from hippt.design_tokens import VIEWPORT_W  # noqa: PLC0415 — inline to survive formatter

    ref_dim = VIEWPORT_W
    px_values = sorted(set(int(v * ref_dim / 100) for v in padding_values if v > 0))

    if not px_values:
        return "0px", []

    if len(px_values) == 1:
        return f"{px_values[0]}px", px_values

    # Find GCD as base rhythm
    base = px_values[0]
    for v in px_values[1:]:
        base = math.gcd(base, v)

    if base < 4:
        base = min(px_values)

    # Generate scale from base
    scale = sorted(set(v for v in px_values if v > 0))
    return f"{base}px", scale


def pptx_to_tokens(pptx_path: str, slug: str | None = None) -> dict:
    """Convert PPTX to design tokens matching reference-extract YAML schema."""
    path = Path(pptx_path)
    if slug is None:
        slug = path.stem.lower().replace(" ", "-")

    analysis = analyze_file(pptx_path)
    summary = analysis["summary"]
    font_families = extract_font_families(pptx_path)

    # --- Palette ---
    raw_palette = summary["color_palette"]
    total_color_count = sum(c["count"] for c in raw_palette)
    palette = []
    for rank, item in enumerate(raw_palette):
        role = classify_color_role(item["hex"], rank, total_color_count, item["count"])
        palette.append(
            {
                "hex": item["hex"],
                "role": role,
                "frequency": item["count"],
            }
        )

    # --- Typography ---
    font_stats = summary.get("font_stats", {})
    all_sizes = []
    for slide in analysis["slides"]:
        all_sizes.extend(slide["font_sizes_pt"])
    all_sizes.sort()

    p50 = all_sizes[len(all_sizes) // 2] if all_sizes else 14
    p75 = all_sizes[int(len(all_sizes) * 0.75)] if all_sizes else 24

    # Aggregate font families by family name
    family_agg: dict[str, dict] = {}
    for f in font_families:
        name = f["family"]
        if name not in family_agg:
            family_agg[name] = {"family": name, "sizes": [], "total_freq": 0}
        family_agg[name]["sizes"].append(float(f["size"].rstrip("pt")))
        family_agg[name]["total_freq"] += f["frequency"]

    typography = []
    for fam in sorted(family_agg.values(), key=lambda x: -x["total_freq"]):
        avg_size = sum(fam["sizes"]) / len(fam["sizes"])
        role = classify_font_role(avg_size)
        typography.append(
            {
                "family": fam["family"],
                "weight": 400,
                "size": f"{avg_size:.0f}pt",
                "role": role,
                "frequency": fam["total_freq"],
            }
        )

    # --- Spacing ---
    pad = summary.get("avg_padding_pct", {})
    pad_values = [v for v in pad.values() if v > 0]
    rhythm, scale = compute_spacing_rhythm(pad_values)

    # --- Hierarchy weight ratio ---
    title_range = summary.get("title_font_pt", [0, 0])
    body_range = summary.get("body_font_pt", [0, 0])
    title_max = title_range[1] if title_range[1] > 0 else p75
    body_avg = (body_range[0] + body_range[1]) / 2 if body_range[1] > 0 else p50
    hierarchy_ratio = round(title_max / body_avg, 2) if body_avg > 0 else 1.0

    tokens = {
        "source_file": path.name,
        "source_type": "pptx",
        "extracted_at": str(date.today()),
        "slide_count": analysis["slide_count"],
        "palette": palette,
        "typography": typography[:10],
        "spacing": {
            "rhythm": rhythm,
            "values": scale,
        },
        "canvas_fill": {
            "avg": summary["avg_canvas_fill"],
            "median": summary["median_canvas_fill"],
            "min": summary["min_canvas_fill"],
            "max": summary["max_canvas_fill"],
        },
        "layout_density": summary.get("layout_distribution", {}),
        "hierarchy_ratio": hierarchy_ratio,
        "title_font_pt": title_range,
        "body_font_pt": body_range,
        "avg_padding_pct": pad,
        "radii": [],
        "shadows": [],
        "gradients": [],
        "motion": [],
        "css_variables": [],
        "element_focus": None,
    }
    return tokens


def main() -> None:
    if len(sys.argv) < 2:
        print("Usage: pptx_to_tokens.py <file.pptx> [--slug NAME]", file=sys.stderr)
        sys.exit(1)

    pptx_path = sys.argv[1]
    slug = None
    if "--slug" in sys.argv:
        idx = sys.argv.index("--slug")
        if idx + 1 < len(sys.argv):
            slug = sys.argv[idx + 1]

    tokens = pptx_to_tokens(pptx_path, slug)

    if slug is None:
        slug = Path(pptx_path).stem.lower().replace(" ", "-")

    out_dir = Path("output/design")
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / f"ref-{slug}.yaml"

    with open(out_path, "w") as f:
        yaml.dump(
            tokens, f, default_flow_style=False, sort_keys=False, allow_unicode=True
        )

    print(f"Design tokens written to {out_path}", file=sys.stderr)
    print(yaml.dump(tokens, default_flow_style=False, sort_keys=False))


if __name__ == "__main__":
    main()
