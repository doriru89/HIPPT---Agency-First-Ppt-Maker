#!/usr/bin/env python3
"""Final PPTX engine — DOM extraction via Playwright, exact coordinate mapping.

Renders approved HTML at 960×540 (matching PPTX 96 DPI native), extracts exact
element positions/styles via getBoundingClientRect + getComputedStyle, builds
editable PPTX with python-pptx.

Counterpart: slides_to_pptx.py = draft PPTX engine (semantic sidecar JSON, fast,
no browser). This file = final PPTX engine (DOM extraction, Playwright, slow but
exact). Both coexist permanently.

Usage:
    uv run python TOOLS/scripts/html_to_pptx.py <html_path> --out <output.pptx> [--debug]

Failure taxonomy:
    FATAL    — raise + exit (Playwright crash, no slides found, save failure)
    RECOVER  — log warning + skip element (zero-size, corrupt base64, parse error)
    DEGRADE  — rasterize instead of native (SVG, complex CSS, JS error in extraction)
"""

from __future__ import annotations

import argparse
import base64
import json
import logging
import re
import sys
import tempfile
import threading
from http.server import HTTPServer, SimpleHTTPRequestHandler
from pathlib import Path


from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from hippt.design_tokens import (
    ALIGN_CSS_TO_PPTX,
    CSS_PX_TO_PT,
    SLIDE_HEIGHT_IN as SLIDE_HEIGHT,
    SLIDE_WIDTH_IN as SLIDE_WIDTH,
    TEXT_PAD_IN,
    VIEWPORT_H,
    VIEWPORT_W,
    parse_css_px,
)

log = logging.getLogger(__name__)

# Render at 960×540 so 1 CSS px = 1/96 inch = native PPTX DPI.
# At 1920×1080, there's a 2x mismatch: positions scale at 0.5x but
# CSS fonts at 0.75 pt/px are 2x too large for the scaled boxes.
# At 960×540, both use the same coordinate system. CSS clamp() values
# naturally produce PPTX-appropriate font sizes (Title ~40pt, Body ~11pt).

ALIGN_MAP = {k: getattr(PP_ALIGN, v) for k, v in ALIGN_CSS_TO_PPTX.items()}

# --- Debug JSON schema ---
# When --debug is set, output/debug/layout-classified.json is written with this shape:
# {
#   "viewport": { "w": 1920, "h": 1080 },
#   "slides": [{
#     "index": 0,
#     "background": "rgb(0, 53, 107)",
#     "bbox": { "x": 0, "y": 0, "w": 1920, "h": 1080 },
#     "elements": [{
#       "type": "text" | "image" | "shape" | "table",
#       "render": "native" | "rasterize" | "skip",
#       "bbox": { "x": number, "y": number, "w": number, "h": number },
#       "style": { fontFamily, fontSize, fontWeight, color, textAlign, backgroundColor,
#                  borderRadius, borderWidth, borderColor, opacity, lineHeight,
#                  paddingLeft, paddingTop, paddingRight, paddingBottom },
#       "runs": [{ "text": str, "style": { fontWeight, fontStyle, color, fontSize } }],
#       "text": str,  // flattened full text
#       "overflow": bool,
#       "src": str,   // image only
#       "rows": [[{ "text": str, "style": dict, "colspan": int, "rowspan": int }]],
#       "isDecorative": bool,
#       "selector": str  // CSS selector for rasterize targeting
#     }]
#   }]
# }


# ============================================================
# Phase 0: Coordinate mapping & color parsing
# ============================================================


def px_to_in(px: float, viewport_px: float, slide_in: float) -> float:
    return (px / viewport_px) * slide_in


def parse_css_color(css: str) -> tuple[RGBColor, float]:
    """Parse css color string → (RGBColor, alpha 0.0-1.0).

    Handles: rgb(r,g,b), rgba(r,g,b,a), #RRGGBB, #RGB, transparent,
    linear-gradient/radial-gradient (extracts first hex as solid fallback).
    Returns (RGBColor(0,0,0), 0.0) for transparent/unparseable.
    """
    from hippt.design_tokens import parse_css_color as _dt_parse_color

    r, g, b, a = _dt_parse_color(css)
    return RGBColor(r, g, b), a


_parse_px = parse_css_px


_CSS_ANGLE_MAP = {
    "to top": 0.0,
    "to top right": 45.0,
    "to right": 90.0,
    "to bottom right": 135.0,
    "to bottom": 180.0,
    "to bottom left": 225.0,
    "to left": 270.0,
    "to top left": 315.0,
}


def _split_gradient_args(body: str) -> list[str]:
    """Split gradient arguments respecting parentheses in rgb()/rgba()."""
    parts: list[str] = []
    depth = 0
    current: list[str] = []
    for ch in body:
        if ch == "(":
            depth += 1
            current.append(ch)
        elif ch == ")":
            depth -= 1
            current.append(ch)
        elif ch == "," and depth == 0:
            parts.append("".join(current).strip())
            current = []
        else:
            current.append(ch)
    if current:
        parts.append("".join(current).strip())
    return parts


def parse_css_gradient(css: str) -> dict | None:
    """Parse CSS linear-gradient → {angle, stops: [(RGBColor, position), ...]}."""
    if not css or "gradient" not in css:
        return None
    m = re.search(r"linear-gradient\(", css, re.IGNORECASE)
    if not m:
        return None
    start = m.end()
    depth = 1
    end = start
    while end < len(css) and depth > 0:
        if css[end] == "(":
            depth += 1
        elif css[end] == ")":
            depth -= 1
        end += 1
    body = css[start : end - 1].strip()

    parts = _split_gradient_args(body)
    if not parts:
        return None

    angle = 180.0  # default: top-to-bottom
    color_parts = parts
    first = parts[0].strip()
    if first in _CSS_ANGLE_MAP:
        angle = _CSS_ANGLE_MAP[first]
        color_parts = parts[1:]
    elif first.endswith("deg"):
        try:
            angle = float(first.replace("deg", ""))
        except ValueError:
            pass
        color_parts = parts[1:]

    stops = []
    for i, part in enumerate(color_parts):
        color_str = re.sub(r"\s+[\d.]+%\s*$", "", part).strip()
        pos_m = re.search(r"([\d.]+)%", part)
        pos = (
            float(pos_m.group(1)) / 100.0
            if pos_m
            else (i / max(len(color_parts) - 1, 1))
        )
        pos = max(0.0, min(1.0, pos))
        color, alpha = parse_css_color(color_str)
        if alpha > 0.01:
            stops.append((color, pos))

    if len(stops) < 2:
        return None

    # CSS: 0deg=to-top, 90deg=to-right, 180deg=to-bottom
    # python-pptx: 0=left-to-right, 90=bottom-to-top, 270=top-to-bottom
    pptx_angle = (360.0 - angle + 90.0) % 360.0

    # Preserve first+last stops when >2 (python-pptx default has exactly 2)
    if len(stops) > 2:
        stops = [stops[0], stops[-1]]

    return {"angle": pptx_angle, "stops": stops}


# ============================================================
# Reference PPTX template support
# ============================================================

from hippt.design_tokens import HEADING_MIN_PT as _HEADING_SIZE_THRESHOLD_PT  # noqa: E402, PLC0415

_TEXT_PAD_IN = TEXT_PAD_IN  # E-PPTX-005: fixed constant, never proportional


def _prepare_template(reference_path: str) -> "Presentation":
    """Load reference PPTX, validate, remove all slides. Returns clean template."""
    ref = Path(reference_path)
    if not ref.exists():
        log.error("FATAL: Reference PPTX not found: %s", ref)
        sys.exit(1)

    try:
        prs = Presentation(str(ref))
    except Exception as e:
        log.error("FATAL: Cannot open reference PPTX (corrupt?): %s", e)
        sys.exit(1)

    ar = prs.slide_width / prs.slide_height
    if abs(ar - 16 / 9) > 0.05:
        actual_w = prs.slide_width / 914400
        actual_h = prs.slide_height / 914400
        log.error(
            "FATAL: Reference PPTX aspect ratio %.2f:1 does not match required 16:9 "
            '(slide dimensions: %.1f"x%.1f")',
            ar,
            actual_w,
            actual_h,
        )
        sys.exit(1)

    layout_count = sum(1 for _ in prs.slide_layouts)
    if layout_count == 0:
        log.error("FATAL: Reference PPTX has no slide layouts")
        sys.exit(1)

    # Strip existing slides via XML manipulation (python-pptx has no public delete API)
    sldIdLst = prs.slides._sldIdLst
    ns_r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    slide_count = len(sldIdLst)
    while len(sldIdLst) > 0:
        rId = sldIdLst[0].get(f"{ns_r}id")
        if rId:
            try:
                prs.part.drop_rel(rId)
            except KeyError:
                log.warning(
                    "Relationship %s not found during slide strip — skipping", rId
                )
        del sldIdLst[0]

    if len(prs.slides._sldIdLst) != 0:
        log.error("FATAL: Slide stripping failed — %d slides remain", len(prs.slides))
        sys.exit(1)

    log.info(
        "Reference template: %s (%d layouts, %d slides stripped)",
        ref.name,
        layout_count,
        slide_count,
    )
    return prs


def _find_blank_layout(prs: "Presentation"):
    """Find best layout — 'Blank' by name, or fewest placeholders as fallback."""
    for layout in prs.slide_layouts:
        if layout.name.lower().strip() == "blank":
            log.info("Using layout '%s' (exact match)", layout.name)
            return layout

    for layout in prs.slide_layouts:
        if "blank" in layout.name.lower():
            log.info("Using layout '%s' (partial match)", layout.name)
            return layout

    best = None
    best_count = float("inf")
    for layout in prs.slide_layouts:
        ph_count = len(layout.placeholders)
        if ph_count < best_count:
            best_count = ph_count
            best = layout

    if best is None:
        names = ", ".join(f"'{la.name}'" for la in prs.slide_layouts)
        log.error("FATAL: No usable layout in reference PPTX. Available: %s", names)
        sys.exit(1)

    log.warning(
        "No 'Blank' layout found. Using '%s' (%d placeholders)",
        best.name,
        best_count,
    )
    return best


def _extract_theme_fonts(prs: "Presentation") -> dict[str, str | None]:
    """Extract heading and body font families from reference PPTX theme XML."""
    fonts: dict[str, str | None] = {"heading": None, "body": None}
    try:
        from lxml import (
            etree,
        )  # Transitive dep via python-pptx; lazy to avoid hard requirement

        if not prs.slide_masters:
            log.warning("No slide masters in reference PPTX — theme fonts unavailable")
            return fonts

        master = prs.slide_masters[0]
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        for rel in master.part.rels.values():
            if "theme" in rel.reltype:
                theme_xml = etree.fromstring(rel.target_part.blob)
                major = theme_xml.find(".//a:majorFont/a:latin", ns)
                minor = theme_xml.find(".//a:minorFont/a:latin", ns)
                if major is not None:
                    fonts["heading"] = major.get("typeface")
                if minor is not None:
                    fonts["body"] = minor.get("typeface")
                break
    except Exception as e:
        log.warning("Could not extract theme fonts: %s", e)

    if fonts["heading"] or fonts["body"]:
        log.info(
            "Theme fonts: heading='%s', body='%s'", fonts["heading"], fonts["body"]
        )
    else:
        log.warning("No theme fonts found — HTML font families will be preserved")
    return fonts


# ============================================================
# Phase 1: Extract — Playwright DOM walk
# ============================================================

_EXTRACT_JS = """
() => {
    const slide = document.querySelector('.slide.active');
    if (!slide) return null;

    const slideRect = slide.getBoundingClientRect();
    const cs = getComputedStyle(slide);

    function getStyle(el) {
        const s = getComputedStyle(el);
        return {
            fontFamily: s.fontFamily,
            fontSize: s.fontSize,
            fontWeight: s.fontWeight,
            fontStyle: s.fontStyle,
            color: s.color,
            textAlign: s.textAlign,
            backgroundColor: s.backgroundColor,
            borderRadius: s.borderRadius,
            borderWidth: s.borderTopWidth,
            borderColor: s.borderTopColor,
            opacity: s.opacity,
            lineHeight: s.lineHeight,
            paddingLeft: s.paddingLeft,
            paddingTop: s.paddingTop,
            paddingRight: s.paddingRight,
            paddingBottom: s.paddingBottom,
            textDecoration: s.textDecorationLine,
        };
    }

    function getRuns(el) {
        const runs = [];
        const walker = document.createTreeWalker(el, NodeFilter.SHOW_TEXT);
        let node;
        while (node = walker.nextNode()) {
            const text = node.textContent;
            if (!text || !text.trim()) continue;
            const parent = node.parentElement;
            const ps = getComputedStyle(parent);
            runs.push({
                text: text,
                style: {
                    fontWeight: ps.fontWeight,
                    fontStyle: ps.fontStyle,
                    color: ps.color,
                    fontSize: ps.fontSize,
                    fontFamily: ps.fontFamily,
                    textDecoration: ps.textDecorationLine,
                }
            });
        }
        return runs;
    }

    function extractElement(el) {
        const rect = el.getBoundingClientRect();
        const w = rect.width, h = rect.height;
        if (w < 2 || h < 2) return null;

        const style = getStyle(el);
        const tag = el.tagName.toLowerCase();
        const bbox = {
            x: rect.left - slideRect.left,
            y: rect.top - slideRect.top,
            w: w, h: h
        };

        // Image
        if (tag === 'img') {
            return {
                type: 'image', bbox, style,
                src: el.src,
                naturalWidth: el.naturalWidth,
                naturalHeight: el.naturalHeight,
                selector: buildSelector(el)
            };
        }

        // Table
        if (tag === 'table') {
            const rows = [];
            el.querySelectorAll('tr').forEach(tr => {
                const cells = [];
                tr.querySelectorAll('td, th').forEach(cell => {
                    cells.push({
                        text: cell.textContent.trim(),
                        style: getStyle(cell),
                        colspan: cell.colSpan || 1,
                        rowspan: cell.rowSpan || 1,
                        isHeader: cell.tagName === 'TH',
                    });
                });
                if (cells.length) rows.push(cells);
            });
            return { type: 'table', bbox, style, rows, selector: buildSelector(el) };
        }

        // SVG
        if (tag === 'svg' || el.querySelector('svg')) {
            return { type: 'image', bbox, style, src: null, selector: buildSelector(el) };
        }

        const text = el.textContent.trim();
        const bg = style.backgroundColor;
        const hasBg = bg && bg !== 'rgba(0, 0, 0, 0)' && bg !== 'transparent';

        // Accent bar (thin element with background — must precede general decorative)
        if (hasBg && h < 6 && !text) {
            return {
                type: 'shape', bbox, style,
                isDecorative: true,
                selector: buildSelector(el)
            };
        }

        // Decorative shape (background + no meaningful text)
        if (hasBg && !text) {
            return {
                type: 'shape', bbox, style,
                isDecorative: true,
                selector: buildSelector(el)
            };
        }

        // Text element
        if (text) {
            const runs = getRuns(el);
            const overflow = el.scrollHeight > el.clientHeight + 2;
            return {
                type: 'text', bbox, style,
                text: text,
                runs: runs.length > 0 ? runs : [{ text, style: getStyle(el) }],
                overflow,
                isDecorative: false,
                selector: buildSelector(el)
            };
        }

        return null;
    }

    function buildSelector(el) {
        if (el.id) return '#' + el.id;
        const classes = Array.from(el.classList).join('.');
        const tag = el.tagName.toLowerCase();
        return classes ? tag + '.' + classes : tag;
    }

    // Walk direct children of key containers
    const elements = [];
    const seen = new Set();

    // E-PRES-037: check ANY visible descendant, not just direct children.
    // Catches wrapper divs with zero dimensions hiding visible grandchildren.
    function hasVisibleDescendant(el) {
        for (const c of el.children) {
            const cr = c.getBoundingClientRect();
            if (cr.width > 2 && cr.height > 2) return true;
            if (hasVisibleDescendant(c)) return true;
        }
        return false;
    }

    function tryAdd(el) {
        const extracted = extractElement(el);
        if (extracted) {
            const key = `${Math.round(extracted.bbox.x)},${Math.round(extracted.bbox.y)},${Math.round(extracted.bbox.w)},${Math.round(extracted.bbox.h)}`;
            if (!seen.has(key)) {
                seen.add(key);
                elements.push(extracted);
            }
        }
    }

    function walk(container, depth) {
        if (depth > 6) return;
        for (const child of container.children) {
            const tag = child.tagName.toLowerCase();
            if (['script', 'style', 'link', 'meta', 'noscript'].includes(tag)) continue;

            const rect = child.getBoundingClientRect();
            if (rect.width < 2 || rect.height < 2) continue;

            const hasDirectText = Array.from(child.childNodes)
                .some(n => n.nodeType === 3 && n.textContent.trim());
            const isLeaf = child.children.length === 0;
            const isTable = tag === 'table';
            const isImg = tag === 'img';
            const hasBg = (() => {
                const bg = getComputedStyle(child).backgroundColor;
                return bg && bg !== 'rgba(0, 0, 0, 0)' && bg !== 'transparent';
            })();
            const isSmall = rect.height < 6;

            const hasExtractableChildren = !isLeaf && !isTable && !isImg
                && hasVisibleDescendant(child);

            if (isTable || isImg || (hasBg && isSmall) || (isLeaf && !hasExtractableChildren)) {
                tryAdd(child);
            } else if (hasDirectText && !hasExtractableChildren) {
                tryAdd(child);
            }

            if (!isLeaf && !isTable && !isImg) {
                walk(child, depth + 1);
            }
        }
    }

    walk(slide, 0);

    // E-PRES-037: post-extraction containment dedup.
    // If element A spatially contains B and A's text is a superset of B's text,
    // drop A (keep the child, remove the parent container).
    const filtered = elements.filter((a, i) => {
        return !elements.some((b, j) => {
            if (i === j || !a.text || !b.text) return false;
            return a.bbox.x <= b.bbox.x + 2
                && a.bbox.y <= b.bbox.y + 2
                && a.bbox.x + a.bbox.w >= b.bbox.x + b.bbox.w - 2
                && a.bbox.y + a.bbox.h >= b.bbox.y + b.bbox.h - 2
                && a.bbox.h > b.bbox.h * 1.3
                && a.text.includes(b.text);
        });
    });

    return {
        viewport: { w: window.innerWidth, h: window.innerHeight },
        slide: {
            bbox: { x: 0, y: 0, w: slideRect.width, h: slideRect.height },
            background: cs.backgroundColor,
            backgroundImage: cs.backgroundImage !== 'none' ? cs.backgroundImage : null,
            slideType: slide.dataset.slideType || null,
            elements: filtered
        }
    };
}
"""


def _extract_slides(html_path: str) -> list[dict]:
    """Launch Playwright, serve HTML, extract all slides."""
    html_dir = str(Path(html_path).parent)
    html_name = Path(html_path).name

    class QuietHandler(SimpleHTTPRequestHandler):
        def __init__(self, *args, **kwargs):
            super().__init__(*args, directory=html_dir, **kwargs)

        def log_message(self, *args):
            pass

    server = HTTPServer(("127.0.0.1", 0), QuietHandler)
    port = server.server_address[1]
    thread = threading.Thread(target=server.serve_forever, daemon=True)
    thread.start()

    slides_data = []
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(
                viewport={"width": VIEWPORT_W, "height": VIEWPORT_H}
            )
            page.goto(f"http://127.0.0.1:{port}/{html_name}", wait_until="networkidle")
            page.wait_for_timeout(500)

            slide_count = page.evaluate("document.querySelectorAll('.slide').length")
            if slide_count == 0:
                log.error("FATAL: no .slide elements found in HTML")
                sys.exit(1)

            log.info("Found %d slides", slide_count)

            for i in range(slide_count):
                page.evaluate(f"""
                    document.querySelectorAll('.slide').forEach((s, idx) => {{
                        s.classList.toggle('active', idx === {i});
                    }});
                """)
                page.wait_for_timeout(200)

                data = page.evaluate(_EXTRACT_JS)
                if data and data.get("slide"):
                    data["slide"]["index"] = i
                    # E-PPTX-004: capture rasterize candidates while browser is open
                    for elem in data["slide"].get("elements", []):
                        style = elem.get("style", {})
                        br = style.get("borderRadius", "0")
                        bg = style.get("backgroundColor", "")
                        needs_capture = (
                            "50%" in br or "100%" in br or "gradient" in (bg or "")
                        )
                        sel = elem.get("selector")
                        if needs_capture and sel:
                            try:
                                el_handle = page.query_selector(f".slide.active {sel}")
                                if el_handle:
                                    screenshot = el_handle.screenshot(type="png")
                                    elem["rasterize_data"] = base64.b64encode(
                                        screenshot
                                    ).decode()
                            except Exception as e:
                                log.warning(
                                    "RECOVER: rasterize capture failed for %s: %s",
                                    sel,
                                    e,
                                )
                    slides_data.append(data)
                    elem_count = len(data["slide"].get("elements", []))
                    log.info("  Slide %d: %d elements extracted", i + 1, elem_count)
                else:
                    log.warning("RECOVER: slide %d returned no data, skipping", i + 1)

            browser.close()
    finally:
        server.shutdown()

    return slides_data


# ============================================================
# Phase 2: Classify — native | rasterize | skip
# ============================================================


def _classify_elements(slides_data: list[dict]) -> list[dict]:
    """Tag each element with render classification."""
    for sd in slides_data:
        for elem in sd["slide"].get("elements", []):
            bbox = elem.get("bbox", {})
            w, h = bbox.get("w", 0), bbox.get("h", 0)

            if w < 2 or h < 2:
                elem["render"] = "skip"
                continue

            etype = elem.get("type", "")

            if etype == "image" and elem.get("src") is None:
                elem["render"] = "rasterize"
                continue

            style = elem.get("style", {})
            bg = style.get("backgroundColor", "")
            has_gradient = "gradient" in bg if bg else False
            has_complex = has_gradient

            if has_complex:
                elem["render"] = "rasterize"
                continue

            # E-PPTX-004: border-radius > 40% of smaller dimension → rasterize
            border_radius_str = style.get("borderRadius", "0")
            needs_rasterize_br = False
            if "50%" in border_radius_str or "100%" in border_radius_str:
                needs_rasterize_br = True
            elif "/" in border_radius_str:
                needs_rasterize_br = True
            else:
                br_val = _parse_px(border_radius_str)
                min_dim = min(w, h)
                if br_val > 4 and min_dim > 0 and br_val / min_dim > 0.4:
                    needs_rasterize_br = True

            if needs_rasterize_br:
                elem["render"] = "rasterize"
                continue

            if etype in ("text", "shape", "table", "image"):
                elem["render"] = "native"
            else:
                elem["render"] = "skip"

            if elem.get("overflow"):
                log.warning(
                    "RECOVER: text overflow on slide %d — '%s'",
                    sd["slide"].get("index", 0) + 1,
                    (elem.get("text", ""))[:40],
                )

    return slides_data


# ============================================================
# Phase 2.5: Layout quality checks
# Adapted from AeSlides (Pan et al., arXiv 2604.22840, 2026).
# Verifiable metrics outperform VLM scoring for layout flaws.
# Thresholds: TOOLS/lib/design_tokens.py (validated via test_design_tokens.py).
# ============================================================


def _check_layout_quality(slides_data: list[dict]) -> None:
    """Warn on element collisions and low content fill per slide."""
    for sd in slides_data:
        slide_idx = sd["slide"].get("index", 0) + 1
        elems = [
            e
            for e in sd["slide"].get("elements", [])
            if e.get("render") not in ("skip", None)
        ]
        if not elems:
            continue

        viewport_area = VIEWPORT_W * VIEWPORT_H

        # --- Excessive whitespace (AeSlides F1=0.80) ---
        content_area = 0
        for e in elems:
            bbox = e.get("bbox", {})
            content_area += bbox.get("w", 0) * bbox.get("h", 0)
        fill_ratio = content_area / viewport_area if viewport_area else 0
        if fill_ratio < 0.40:
            log.warning(
                "Slide %d: content fill %.0f%% — likely content islands (E-PRES-006)",
                slide_idx,
                fill_ratio * 100,
            )

        # --- Element collision (AeSlides F1=0.83) ---
        bboxes = []
        for e in elems:
            b = e.get("bbox", {})
            x, y, w, h = b.get("x", 0), b.get("y", 0), b.get("w", 0), b.get("h", 0)
            if w > 0 and h > 0:
                bboxes.append((x, y, x + w, y + h, (e.get("text", "") or "")[:30]))
        collisions = 0
        for i in range(len(bboxes)):
            for j in range(i + 1, len(bboxes)):
                ax1, ay1, ax2, ay2, _ = bboxes[i]
                bx1, by1, bx2, by2, _ = bboxes[j]
                ix1, iy1 = max(ax1, bx1), max(ay1, by1)
                ix2, iy2 = min(ax2, bx2), min(ay2, by2)
                if ix1 < ix2 and iy1 < iy2:
                    overlap_area = (ix2 - ix1) * (iy2 - iy1)
                    smaller = min((ax2 - ax1) * (ay2 - ay1), (bx2 - bx1) * (by2 - by1))
                    if smaller > 0 and overlap_area / smaller > 0.5:
                        collisions += 1
        if collisions:
            log.warning(
                "Slide %d: %d element collision(s) detected — text may overlap",
                slide_idx,
                collisions,
            )


# ============================================================
# Phase 3: Build — python-pptx with exact coordinates
# ============================================================


def _build_pptx(
    slides_data: list[dict],
    output_path: str,
    reference_path: str | None = None,
) -> list[int]:
    """Build PPTX from classified extraction data. Returns per-slide element counts."""
    theme_fonts: dict[str, str | None] = {"heading": None, "body": None}

    if reference_path:
        prs = _prepare_template(reference_path)
        theme_fonts = _extract_theme_fonts(prs)
        blank_layout = _find_blank_layout(prs)
    else:
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_WIDTH)
        prs.slide_height = Inches(SLIDE_HEIGHT)
        blank_layout = _find_blank_layout(prs)
    element_counts = []

    for sd in slides_data:
        slide_info = sd["slide"]
        slide = prs.slides.add_slide(blank_layout)
        count = 0

        bg_css = slide_info.get("background", "")
        bg_image_css = slide_info.get("backgroundImage") or ""
        grad = parse_css_gradient(bg_image_css)
        if grad and len(grad["stops"]) >= 2:
            fill = slide.background.fill
            fill.gradient()
            fill.gradient_angle = grad["angle"]
            for i, (color, pos) in enumerate(grad["stops"]):
                fill.gradient_stops[i].color.rgb = color
                fill.gradient_stops[i].position = pos
            log.info(
                "  Slide bg: gradient angle %.0f°",
                grad["angle"],
            )
        else:
            bg_color, bg_alpha = parse_css_color(bg_css)
            if bg_alpha > 0.01:
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = bg_color

        vp = sd.get("viewport", {"w": VIEWPORT_W, "h": VIEWPORT_H})
        vw, vh = vp["w"], vp["h"]

        for elem in slide_info.get("elements", []):
            if elem.get("render") == "skip":
                continue

            etype = elem.get("type", "")
            bbox = elem.get("bbox", {})

            raw_left = px_to_in(bbox["x"], vw, SLIDE_WIDTH)
            raw_top = px_to_in(bbox["y"], vh, SLIDE_HEIGHT)
            raw_w = px_to_in(bbox["w"], vw, SLIDE_WIDTH)
            raw_h = px_to_in(bbox["h"], vh, SLIDE_HEIGHT)

            # Fixed right-side padding for font metric differences.
            # MUST NOT shift left position — proportional padding destroys grid alignment.
            if etype == "text":
                raw_w = min(raw_w + _TEXT_PAD_IN, SLIDE_WIDTH - raw_left)
                raw_h = min(raw_h + _TEXT_PAD_IN, SLIDE_HEIGHT - raw_top)

            left = Inches(raw_left)
            top = Inches(raw_top)
            width = Inches(raw_w)
            height = Inches(raw_h)

            if width.inches < 0.05 or height.inches < 0.05:
                continue

            try:
                if elem.get("render") == "rasterize":
                    raster_b64 = elem.get("rasterize_data")
                    if raster_b64:
                        img_bytes = base64.b64decode(raster_b64)
                        with tempfile.NamedTemporaryFile(
                            suffix=".png", delete=False
                        ) as tmp:
                            tmp.write(img_bytes)
                            tmp_path = tmp.name
                        slide.shapes.add_picture(tmp_path, left, top, width, height)
                        Path(tmp_path).unlink(missing_ok=True)
                    else:
                        log.warning(
                            "RECOVER: rasterize element on slide %d has no data",
                            slide_info.get("index", 0) + 1,
                        )
                    count += 1
                    continue

                if etype == "text":
                    _add_text(slide, elem, left, top, width, height, theme_fonts)
                    count += 1
                elif etype == "image":
                    _add_image(slide, elem, left, top, width, height)
                    count += 1
                elif etype == "shape":
                    _add_shape(slide, elem, left, top, width, height)
                    count += 1
                elif etype == "table":
                    _add_table(slide, elem, left, top, width, height)
                    count += 1
            except Exception as e:
                log.warning(
                    "RECOVER: failed to add %s element on slide %d: %s",
                    etype,
                    slide_info.get("index", 0) + 1,
                    e,
                )

        element_counts.append(count)

    prs.save(output_path)
    log.info("Saved PPTX to %s", output_path)
    return element_counts


def _add_text(slide, elem: dict, left, top, width, height, theme_fonts=None):
    """Add textbox with multi-run support."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)

    style = elem.get("style", {})
    align = ALIGN_MAP.get(style.get("textAlign", "left"), PP_ALIGN.LEFT)

    bg_color_str = style.get("backgroundColor", "")
    bg_color, bg_alpha = parse_css_color(bg_color_str)
    if bg_alpha > 0.01:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    border_radius = _parse_px(style.get("borderRadius", "0"))
    if border_radius > 4:
        txBox.fill.solid()
        if bg_alpha > 0.01:
            txBox.fill.fore_color.rgb = bg_color

    runs = elem.get("runs", [])
    if not runs:
        return

    # Two-pass: flatten runs into paragraph groups, then build.
    # Each \n starts a new paragraph group. No mutable first_para leaking.
    para_groups: list[list[tuple[str, dict]]] = [[]]

    for run_data in runs:
        text = run_data.get("text", "")
        run_style = run_data.get("style", {})
        if not text:
            continue
        if "\n" in text:
            lines = text.split("\n")
            for li, line in enumerate(lines):
                if li > 0:
                    para_groups.append([])
                stripped = line.strip()
                if stripped:
                    para_groups[-1].append((stripped, run_style))
        else:
            para_groups[-1].append((text, run_style))

    for pi, group in enumerate(para_groups):
        if not group:
            continue
        para = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        para.alignment = align
        for text, run_style in group:
            run = para.add_run()
            run.text = text
            _apply_run_style(run, run_style, theme_fonts)


def _apply_run_style(run, style: dict, theme_fonts=None):
    """Apply font styling to a python-pptx run."""
    font = run.font

    font_size = _parse_px(style.get("fontSize", "14px"))
    if font_size > 0:
        font.size = Pt(font_size * CSS_PX_TO_PT)

    _font_applied = False
    if theme_fonts and (theme_fonts.get("heading") or theme_fonts.get("body")):
        pt = font_size * CSS_PX_TO_PT
        if pt >= _HEADING_SIZE_THRESHOLD_PT and theme_fonts.get("heading"):
            font.name = theme_fonts["heading"]
            _font_applied = True
        elif theme_fonts.get("body"):
            font.name = theme_fonts["body"]
            _font_applied = True

    if not _font_applied:
        font_family = style.get("fontFamily", "")
        if font_family:
            clean = font_family.split(",")[0].strip().strip("'\"")
            font.name = clean

    weight = style.get("fontWeight", "400")
    font.bold = weight in ("bold", "700", "800", "900", "600")

    font.italic = style.get("fontStyle") == "italic"

    color_str = style.get("color", "")
    if color_str:
        color, alpha = parse_css_color(color_str)
        if alpha > 0.01:
            font.color.rgb = color

    if style.get("textDecoration") == "line-through":
        font.strikethrough = True


def _add_image(slide, elem: dict, left, top, width, height):
    """Add image from src (base64, file path, or URL)."""
    src = elem.get("src", "")
    if not src:
        log.warning("RECOVER: image with no src, skipping")
        return

    if src.startswith("data:"):
        m = re.match(r"data:[^;]+;base64,(.*)", src)
        if not m:
            log.warning("RECOVER: unparseable data URI, skipping")
            return
        try:
            img_data = base64.b64decode(m.group(1))
        except Exception:
            log.warning("RECOVER: corrupt base64 image data, skipping")
            return
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp.write(img_data)
            tmp_path = tmp.name
        slide.shapes.add_picture(tmp_path, left, top, width, height)
        Path(tmp_path).unlink(missing_ok=True)
    elif src.startswith(("http://", "https://")):
        log.warning(
            "RECOVER: external URL image '%s' — skipping (use local assets)", src[:60]
        )
    elif Path(src).exists():
        slide.shapes.add_picture(src, left, top, width, height)
    else:
        log.warning("RECOVER: image not found '%s', skipping", src[:80])


def _add_shape(slide, elem: dict, left, top, width, height):
    """Add rectangle or rounded rectangle shape."""
    style = elem.get("style", {})
    border_radius = _parse_px(style.get("borderRadius", "0"))

    if border_radius > 4:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

    bg_str = style.get("backgroundColor", "")
    bg_color, bg_alpha = parse_css_color(bg_str)
    if bg_alpha > 0.01:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
    else:
        shape.fill.background()

    border_str = style.get("borderColor", "")
    border_w = _parse_px(style.get("borderWidth", "0"))
    if border_w > 0 and border_str:
        b_color, b_alpha = parse_css_color(border_str)
        if b_alpha > 0.01:
            shape.line.color.rgb = b_color
            shape.line.width = Pt(border_w * CSS_PX_TO_PT)
    else:
        shape.line.fill.background()


def _add_table(slide, elem: dict, left, top, width, height):
    """Add native table with per-cell styling."""
    rows_data = elem.get("rows", [])
    if not rows_data:
        return

    n_rows = len(rows_data)
    n_cols = max(len(r) for r in rows_data) if rows_data else 1

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    for ri, row in enumerate(rows_data):
        for ci, cell_data in enumerate(row):
            if ci >= n_cols:
                break
            cell = table.cell(ri, ci)
            cell.text = cell_data.get("text", "")

            cell_style = cell_data.get("style", {})
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    _apply_run_style(run, cell_style)

            bg_str = cell_style.get("backgroundColor", "")
            if bg_str:
                bg_c, bg_a = parse_css_color(bg_str)
                if bg_a > 0.01:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = bg_c


# ============================================================
# Validation (imported pattern from slides_to_pptx.py)
# ============================================================


def _validate_roundtrip(
    pptx_path: str, expected_slides: int, expected_elements: list[int]
) -> bool:
    """Re-open PPTX and verify slide/element counts."""
    try:
        prs = Presentation(pptx_path)
        actual = len(prs.slides)
        if actual != expected_slides:
            log.error(
                "Round-trip FAIL: expected %d slides, got %d", expected_slides, actual
            )
            return False
        for i, slide in enumerate(prs.slides):
            actual_shapes = len(slide.shapes)
            if i < len(expected_elements) and actual_shapes < expected_elements[i]:
                log.warning(
                    "Slide %d: expected >= %d shapes, got %d",
                    i + 1,
                    expected_elements[i],
                    actual_shapes,
                )
        log.info("Round-trip OK: %d slides verified", actual)
        return True
    except Exception as e:
        log.error("Round-trip validation failed: %s", e)
        return False


# ============================================================
# CLI
# ============================================================


def main():
    parser = argparse.ArgumentParser(description="HTML → PPTX via DOM extraction")
    parser.add_argument("html", help="Path to HTML presentation file")
    parser.add_argument("--out", default=None, help="Output PPTX path")
    parser.add_argument(
        "--debug", action="store_true", help="Dump classified JSON to output/debug/"
    )
    parser.add_argument(
        "--layout-library", metavar="DIR", help="Extract layout YAMLs to DIR"
    )
    parser.add_argument(
        "--reference",
        metavar="PPTX",
        help="Reference PPTX for theme inheritance (fonts, colors, slide masters)",
    )
    parser.add_argument(
        "--use-layouts",
        metavar="DIR",
        help="Layout library directory for intelligent per-slide position remapping",
    )
    args = parser.parse_args()

    logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")

    html_path = Path(args.html).resolve()
    if not html_path.exists():
        log.error("FATAL: HTML file not found: %s", html_path)
        sys.exit(1)

    out_path = args.out or str(html_path.with_suffix(".pptx"))

    log.info("Phase 1: Extracting from %s", html_path.name)
    slides_data = _extract_slides(str(html_path))

    if not slides_data:
        log.error("FATAL: no slides extracted")
        sys.exit(1)

    log.info("Phase 2: Classifying %d slides", len(slides_data))
    slides_data = _classify_elements(slides_data)

    if args.debug:
        debug_dir = Path("output/debug")
        debug_dir.mkdir(parents=True, exist_ok=True)
        debug_path = debug_dir / "layout-classified.json"
        with open(debug_path, "w") as f:
            json.dump(
                {
                    "viewport": slides_data[0].get("viewport", {})
                    if slides_data
                    else {},
                    "slides": [
                        {
                            "index": sd["slide"].get("index", i),
                            "background": sd["slide"].get("background", ""),
                            "bbox": sd["slide"].get("bbox", {}),
                            "elements": sd["slide"].get("elements", []),
                        }
                        for i, sd in enumerate(slides_data)
                    ],
                },
                f,
                indent=2,
            )
        log.info("Debug JSON written to %s", debug_path)

    if args.layout_library:
        from hippt.layout_utils import write_layout_index, write_layout_yaml
        from hippt.extract_layouts import extract_slide_layout

        layout_dir = Path(args.layout_library)
        layout_dir.mkdir(parents=True, exist_ok=True)
        layouts = []
        viewport = {"w": VIEWPORT_W, "h": VIEWPORT_H}
        for i, sd in enumerate(slides_data):
            layout = extract_slide_layout(sd["slide"], i, viewport)
            if layout:
                write_layout_yaml(layout, layout_dir)
                layouts.append(layout)
        if layouts:
            write_layout_index(layouts, layout_dir)
        log.info("Layout library: %d layouts → %s", len(layouts), layout_dir)

    log.info("Phase 2.5: Layout quality checks")
    _check_layout_quality(slides_data)

    if args.use_layouts:
        from hippt.layout_select import (
            _get_selection_config,
            infer_slide_profile,
            load_layout_library,
            match_elements_to_regions,
            remap_elements,
            select_layout,
        )
        from hippt.layout_utils import load_config as _load_layout_cfg

        _layout_cfg = _load_layout_cfg()
        _sel_cfg = _get_selection_config(_layout_cfg)
        library = load_layout_library(args.use_layouts, cfg=_layout_cfg)
        log.info(
            "Phase 7: Layout selection (%d layouts from %s)",
            len(library),
            args.use_layouts,
        )
        selection_results = []
        for i, sd in enumerate(slides_data):
            profile = infer_slide_profile(sd["slide"], i, len(slides_data))
            match = select_layout(profile, library)
            if match:
                regions = match.layout.get("regions", [])
                elems = [
                    e
                    for e in sd["slide"].get("elements", [])
                    if e.get("render") != "skip"
                ]
                matched_pairs, unmatched = match_elements_to_regions(elems, regions)
                total = len(matched_pairs) + len(unmatched)
                ratio = len(matched_pairs) / total if total else 0
                min_ratio = _sel_cfg.get("min_match_ratio", 0.5)
                remapped = False
                if ratio >= min_ratio:
                    match.matched_regions = matched_pairs
                    match.unmatched_elements = unmatched
                    vp = sd.get("viewport", {"w": VIEWPORT_W, "h": VIEWPORT_H})
                    remap_elements(sd, match, vp)
                    remapped = True
                log.info(
                    "  Slide %d: %s (score %.2f, %d/%d matched%s)",
                    i + 1,
                    match.layout["code"],
                    match.score,
                    len(matched_pairs),
                    total,
                    "" if remapped else ", skipped remap — below min_match_ratio",
                )
                selection_results.append(
                    {
                        "slide": i + 1,
                        "layout": match.layout["code"],
                        "score": round(match.score, 4),
                        "layer1": round(match.layer1_score, 4),
                        "layer2": round(match.layer2_score, 4),
                        "shape": round(match.shape_score, 4),
                        "matched": len(matched_pairs),
                        "unmatched": len(unmatched),
                        "match_ratio": round(ratio, 3),
                        "remapped": remapped,
                        "profile_type": profile.slide_type,
                        "profile_density": profile.density,
                    }
                )
            else:
                log.info("  Slide %d: no layout match, using DOM positions", i + 1)
                selection_results.append(
                    {
                        "slide": i + 1,
                        "layout": None,
                        "profile_type": profile.slide_type,
                        "profile_density": profile.density,
                    }
                )

        log.info("Phase 2.5b: Post-remap layout quality checks")
        _check_layout_quality(slides_data)

        if args.debug:
            debug_dir = Path("output/debug")
            debug_dir.mkdir(parents=True, exist_ok=True)
            import json as _json

            sel_path = debug_dir / "layout-selection.json"
            with open(sel_path, "w") as f:
                _json.dump(
                    {
                        "library_size": len(library),
                        "library_dir": str(args.use_layouts),
                        "slides": selection_results,
                    },
                    f,
                    indent=2,
                )
            log.info("Layout selection debug → %s", sel_path)

    log.info("Phase 3: Building PPTX")
    element_counts = _build_pptx(slides_data, out_path, args.reference)

    log.info("Validating round-trip")
    _validate_roundtrip(out_path, len(slides_data), element_counts)

    total_elements = sum(element_counts)
    log.info(
        "Done: %d slides, %d elements → %s", len(slides_data), total_elements, out_path
    )


if __name__ == "__main__":
    main()
